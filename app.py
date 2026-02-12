#!/usr/bin/env python3
import re
import shutil
import tempfile
import zipfile
import subprocess
import sys
import threading
import uuid
from pathlib import Path
from datetime import datetime
from typing import Optional

from flask import Flask, render_template, request, send_file, jsonify
from werkzeug.utils import secure_filename

from pptx import Presentation
from pptx.util import Emu
from PIL import Image

APP_ROOT = Path(__file__).resolve().parent
REPO_ROOT = APP_ROOT.parent
BEAMER_SCRIPT = REPO_ROOT / "Presentation" / "beamer_to_pptx.py"
VIDEO_SCRIPT = APP_ROOT / "generate_video.py"

ALLOWED_EXT = {".zip"}
MAX_CONTENT_LENGTH = 200 * 1024 * 1024  # 200 MB

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = MAX_CONTENT_LENGTH

BEAMER_RE = re.compile(r"\\documentclass(\[[^\]]*\])?\{beamer\}")
ALLOWED_VIDEO_RATES = {"-20%", "-10%", "+0%", "+10%", "+20%", "+30%", "+40%", "+50%"}
JOBS: dict[str, dict] = {}
JOBS_LOCK = threading.Lock()


def _append_log(job_id: str, message: str) -> None:
    with JOBS_LOCK:
        if job_id in JOBS:
            JOBS[job_id]["log"].append(message)


def _set_job_state(job_id: str, **kwargs) -> None:
    with JOBS_LOCK:
        if job_id in JOBS:
            JOBS[job_id].update(kwargs)


def _safe_extract_zip(zip_path: Path, dest: Path) -> None:
    with zipfile.ZipFile(zip_path, "r") as zf:
        for member in zf.infolist():
            name = member.filename
            if name.startswith("/") or ".." in Path(name).parts:
                raise ValueError("Zip contains unsafe paths")
        zf.extractall(dest)


def _find_main_tex(root: Path) -> Path:
    tex_files = list(root.rglob("*.tex"))
    if not tex_files:
        raise FileNotFoundError("No .tex files found in the zip.")

    # Prefer beamer class
    for tex in tex_files:
        try:
            content = tex.read_text(errors="ignore")
        except Exception:
            continue
        if BEAMER_RE.search(content):
            return tex

    # Fallback preference by common names
    for name in ("presentation.tex", "main.tex", "slides.tex"):
        for tex in tex_files:
            if tex.name.lower() == name:
                return tex

    return tex_files[0]


def _require_tool(binary: str) -> None:
    if shutil.which(binary) is None:
        raise RuntimeError(f"Required tool '{binary}' not found in PATH.")


def _run_cmd(cmd: list[str], cwd: Path, log_cb=None) -> None:
    if log_cb:
        log_cb(f"$ {' '.join(cmd)}")
    proc = subprocess.Popen(
        cmd,
        cwd=cwd,
        stdout=subprocess.PIPE,
        stderr=subprocess.STDOUT,
        text=True,
        bufsize=1,
    )
    assert proc.stdout is not None
    for line in proc.stdout:
        line = line.rstrip()
        if line and log_cb:
            log_cb(line)
    rc = proc.wait()
    if rc != 0:
        raise subprocess.CalledProcessError(rc, cmd)


def _run_beamer_to_pptx(tex_path: Path, out_pptx: Path, log_cb=None) -> None:
    if not BEAMER_SCRIPT.exists():
        raise FileNotFoundError(f"Missing script: {BEAMER_SCRIPT}")

    cmd = [
        sys.executable,
        str(BEAMER_SCRIPT),
        str(tex_path),
        "-o",
        str(out_pptx),
    ]
    _run_cmd(cmd, cwd=tex_path.parent, log_cb=log_cb)


def _run_beamer_to_video(tex_path: Path, out_video: Path, rate: str, log_cb=None) -> None:
    if not VIDEO_SCRIPT.exists():
        raise FileNotFoundError(f"Missing script: {VIDEO_SCRIPT}")

    cmd = [
        sys.executable,
        str(VIDEO_SCRIPT),
        "--tex",
        str(tex_path),
        "--output",
        str(out_video),
        "--rate",
        rate,
    ]
    _run_cmd(cmd, cwd=tex_path.parent, log_cb=log_cb)


def _compile_beamer_to_pdf(tex_path: Path, out_pdf: Path, log_cb=None) -> None:
    _require_tool("latexmk")
    cmd = [
        "latexmk",
        "-pdf",
        "-interaction=nonstopmode",
        "-halt-on-error",
        tex_path.name,
    ]
    _run_cmd(cmd, cwd=tex_path.parent, log_cb=log_cb)
    built_pdf = tex_path.with_suffix(".pdf")
    if not built_pdf.exists():
        raise RuntimeError("latexmk finished but PDF was not created.")
    built_pdf.replace(out_pdf)


def _pdf_to_pngs(pdf_path: Path, out_dir: Path, dpi: int = 300, log_cb=None) -> list[Path]:
    _require_tool("pdftoppm")
    out_dir.mkdir(parents=True, exist_ok=True)
    prefix = out_dir / "slide"
    cmd = [
        "pdftoppm",
        "-r",
        str(dpi),
        "-png",
        str(pdf_path),
        str(prefix),
    ]
    _run_cmd(cmd, cwd=out_dir.parent, log_cb=log_cb)
    images = sorted(out_dir.glob("slide-*.png"))
    if not images:
        raise RuntimeError("No PNGs generated from PDF.")
    return images


def _extract_braced_content(text: str, start: int) -> tuple[str, int]:
    if start >= len(text) or text[start] != "{":
        return "", start
    depth = 0
    i = start
    while i < len(text):
        ch = text[i]
        if ch == "\\" and i + 1 < len(text):
            i += 2
            continue
        if ch == "{":
            depth += 1
        elif ch == "}":
            depth -= 1
            if depth == 0:
                return text[start + 1:i], i + 1
        i += 1
    return text[start + 1 :], len(text)


def _strip_latex_comments(text: str) -> str:
    return re.sub(r"(?<!\\)%.*", "", text)


def _latex_to_plain_text(text: str) -> str:
    text = _strip_latex_comments(text)
    text = re.sub(r"\\begin\{[^}]+\}", "", text)
    text = re.sub(r"\\end\{[^}]+\}", "", text)
    text = text.replace("\\item", "\n- ")
    text = re.sub(r"\\\\", "\n", text)
    text = re.sub(r"\\(?:small|footnotesize|scriptsize|tiny|normalsize|large|Large|huge|Huge)\b", "", text)
    text = re.sub(r"\\textcolor\{[^{}]*\}\{([^{}]*)\}", r"\1", text)
    text = re.sub(r"\\(?:textbf|textit|emph|underline|texttt|textrm|alert|structure)\{([^{}]*)\}", r"\1", text)
    text = re.sub(r"\\[a-zA-Z*]+(?:\[[^\]]*\])?(?:\{[^{}]*\})?", "", text)
    text = text.replace("{", "").replace("}", "")
    text = re.sub(r"\s+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]{2,}", " ", text)
    return text.strip()


def _extract_frame_notes(tex_path: Path) -> list[str]:
    source = tex_path.read_text(errors="ignore")
    source = _strip_latex_comments(source)
    frame_re = re.compile(
        r"\\begin\{frame\}(?:\[[^\]]*\])?(?:\{(?P<title>[^{}]*)\})?(?:\{[^{}]*\})?(?P<body>.*?)\\end\{frame\}",
        re.DOTALL,
    )
    notes = []
    for m in frame_re.finditer(source):
        body = m.group("body")
        note_chunks = []
        idx = 0
        while True:
            note_match = re.search(r"\\narration(?:\[[^\]]*\])?(?:<[^>]*>)?\{", body[idx:])
            if not note_match:
                break
            note_start = idx + note_match.start()
            brace_pos = note_start + note_match.group(0).rfind("{")
            content, next_pos = _extract_braced_content(body, brace_pos)
            cleaned = _latex_to_plain_text(content)
            if cleaned:
                note_chunks.append(cleaned)
            idx = next_pos

        note_text = "\n\n".join(note_chunks).strip()
        notes.append(note_text.strip())
    return notes


def _images_to_pptx(images: list[Path], out_pptx: Path, notes: Optional[list[str]] = None) -> None:
    prs = Presentation()
    # Match slide size to first image aspect ratio
    with Image.open(images[0]) as img:
        w_px, h_px = img.size
    slide_width = 13.333  # inches
    slide_height = slide_width * (h_px / w_px)
    prs.slide_width = Emu(slide_width * 914400)
    prs.slide_height = Emu(slide_height * 914400)

    blank = prs.slide_layouts[6]
    for idx, img_path in enumerate(images):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(img_path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
        if notes and idx < len(notes) and notes[idx]:
            notes_slide = slide.notes_slide
            notes_frame = notes_slide.notes_text_frame
            notes_frame.clear()
            notes_frame.text = notes[idx]
    prs.save(out_pptx)


def _process_job(job_id: str, filename: str, mode: str, payload: bytes, video_rate: str) -> None:
    tmp_dir = Path(tempfile.mkdtemp(prefix="beamer_job_"))
    _set_job_state(job_id, temp_dir=str(tmp_dir))
    log = lambda m: _append_log(job_id, m)

    try:
        zip_path = tmp_dir / filename
        zip_path.write_bytes(payload)
        log(f"Saved upload: {filename}")

        src_dir = tmp_dir / "src"
        src_dir.mkdir(parents=True, exist_ok=True)
        _safe_extract_zip(zip_path, src_dir)
        log("Extracted zip.")

        tex_path = _find_main_tex(src_dir)
        log(f"Using TeX file: {tex_path.relative_to(src_dir)}")

        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        if mode == "editable":
            out_pptx = tmp_dir / f"beamer_editable_{stamp}.pptx"
            _run_beamer_to_pptx(tex_path, out_pptx, log_cb=log)
            _set_job_state(job_id, status="done", output_path=str(out_pptx), output_name=out_pptx.name)
            log("Editable PPTX generated.")
            return

        if mode == "video":
            out_video = tmp_dir / f"beamer_video_{stamp}.mp4"
            _run_beamer_to_video(tex_path, out_video, rate=video_rate, log_cb=log)
            _set_job_state(job_id, status="done", output_path=str(out_video), output_name=out_video.name)
            log("Video generated.")
            return

        out_pdf = tmp_dir / f"beamer_{stamp}.pdf"
        _compile_beamer_to_pdf(tex_path, out_pdf, log_cb=log)
        images_dir = tmp_dir / "images"
        images = _pdf_to_pngs(out_pdf, images_dir, log_cb=log)
        out_pptx = tmp_dir / f"beamer_images_{stamp}.pptx"
        notes = _extract_frame_notes(tex_path)
        _images_to_pptx(images, out_pptx, notes=notes)
        _set_job_state(job_id, status="done", output_path=str(out_pptx), output_name=out_pptx.name)
        log("Image-based PPTX generated.")
    except subprocess.CalledProcessError as e:
        _set_job_state(job_id, status="error", error=f"Command failed: {' '.join(e.cmd)}")
        log(f"Error: command failed with exit code {e.returncode}")
    except Exception as e:
        _set_job_state(job_id, status="error", error=str(e))
        log(f"Error: {e}")


@app.route("/")
def index():
    return render_template("index.html")


@app.route("/convert", methods=["POST"])
def convert():
    if "file" not in request.files:
        return jsonify({"error": "Missing file upload."}), 400

    f = request.files["file"]
    if f.filename == "":
        return jsonify({"error": "No file selected."}), 400

    filename = secure_filename(f.filename)
    ext = Path(filename).suffix.lower()
    if ext not in ALLOWED_EXT:
        return jsonify({"error": "Only .zip files are supported."}), 400

    mode = request.form.get("mode", "editable")
    if mode not in {"editable", "images", "video"}:
        return jsonify({"error": "Invalid mode."}), 400
    video_rate = request.form.get("video_rate", "+20%")
    if video_rate not in ALLOWED_VIDEO_RATES:
        return jsonify({"error": "Invalid narration speed."}), 400

    payload = f.read()
    job_id = uuid.uuid4().hex
    with JOBS_LOCK:
        JOBS[job_id] = {
            "status": "running",
            "log": [f"Started job {job_id} ({mode})"],
            "error": None,
            "output_path": None,
            "output_name": None,
            "temp_dir": None,
        }
    thread = threading.Thread(target=_process_job, args=(job_id, filename, mode, payload, video_rate), daemon=True)
    thread.start()
    return jsonify({"job_id": job_id})


@app.route("/status/<job_id>", methods=["GET"])
def status(job_id: str):
    with JOBS_LOCK:
        job = JOBS.get(job_id)
        if not job:
            return jsonify({"error": "Job not found."}), 404
        return jsonify(
            {
                "status": job["status"],
                "log": "\n".join(job["log"]),
                "error": job["error"],
                "ready": bool(job["output_path"]),
            }
        )


@app.route("/download/<job_id>", methods=["GET"])
def download(job_id: str):
    with JOBS_LOCK:
        job = JOBS.get(job_id)
        if not job:
            return jsonify({"error": "Job not found."}), 404
        output_path = job["output_path"]
        output_name = job["output_name"]

    if not output_path or not output_name:
        return jsonify({"error": "Output not ready."}), 409

    output_file = Path(output_path)
    if not output_file.exists():
        return jsonify({"error": "Output file missing."}), 410

    return send_file(output_file, as_attachment=True, download_name=output_name)


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=8000, debug=True)
