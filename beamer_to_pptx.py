#!/usr/bin/env python3
"""
Beamer LaTeX to PPTX converter.

Parses a Beamer .tex file and generates an editable PPTX with matching
colors, blocks, tikzpicture diagrams, frames, columns, tables, math, etc.

Usage:
    pip install python-pptx matplotlib Pillow qrcode
    python beamer_to_pptx.py presentation.tex -o presentation_parsed.pptx
"""

import argparse
import logging
import re
import subprocess
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from PIL import Image

logger = logging.getLogger("beamer_to_pptx")

# ════════════════════════════════════════════════════════════════════
# 1. DATA CLASSES
# ════════════════════════════════════════════════════════════════════

@dataclass
class TextRun:
    """A span of formatted text."""
    text: str
    bold: bool = False
    italic: bool = False
    color: Optional[str] = None      # color name e.g. "ubblue", "red"
    monospace: bool = False
    size: Optional[str] = None       # "small", "scriptsize", etc.
    href: Optional[str] = None


@dataclass
class TableCell:
    """A single table cell."""
    runs: list = field(default_factory=list)   # list of TextRun
    cellcolor: Optional[str] = None
    multirow: Optional[int] = None
    placeholder: bool = False


@dataclass
class TableRow:
    """A single table row."""
    cells: list = field(default_factory=list)   # list of TableCell
    rowcolor: Optional[str] = None
    is_rule: bool = False


@dataclass
class TexNode:
    """AST node representing a LaTeX element."""
    kind: str
    attrs: dict = field(default_factory=dict)
    children: list = field(default_factory=list)


@dataclass
class PreambleData:
    """Data extracted from the document preamble."""
    colors: dict = field(default_factory=dict)         # name -> (R, G, B)
    beamer_colors: dict = field(default_factory=dict)   # element -> {fg, bg}
    title: str = ""
    subtitle: str = ""
    author: str = ""
    institute: str = ""
    date: str = ""
    aspect_ratio: str = "169"


# ════════════════════════════════════════════════════════════════════
# 2. TEX TOKENIZER
# ════════════════════════════════════════════════════════════════════

class TexTokenizer:
    """Low-level LaTeX text processing: brace matching, command extraction."""

    def __init__(self, source: str):
        self.source = source

    def find_matching_brace(self, pos: int) -> int:
        """Given pos of '{', return pos of matching '}'. Handles nesting."""
        assert self.source[pos] == "{", f"Expected '{{' at pos {pos}, got '{self.source[pos]}'"
        depth = 0
        i = pos
        n = len(self.source)
        while i < n:
            ch = self.source[i]
            if ch == "\\" and i + 1 < n:
                i += 2
                continue
            if ch == "{":
                depth += 1
            elif ch == "}":
                depth -= 1
                if depth == 0:
                    return i
            i += 1
        raise ValueError(f"Unmatched brace at position {pos}")

    def find_matching_bracket(self, pos: int) -> int:
        """Given pos of '[', return pos of matching ']'."""
        assert self.source[pos] == "["
        depth = 0
        i = pos
        n = len(self.source)
        while i < n:
            ch = self.source[i]
            if ch == "\\" and i + 1 < n:
                i += 2
                continue
            if ch == "[":
                depth += 1
            elif ch == "]":
                depth -= 1
                if depth == 0:
                    return i
            i += 1
        raise ValueError(f"Unmatched bracket at position {pos}")

    def find_env_end(self, text: str, env_name: str, start: int = 0) -> int:
        """Find \\end{env_name} matching a \\begin{env_name} at nesting level 0."""
        begin_tag = f"\\begin{{{env_name}}}"
        end_tag = f"\\end{{{env_name}}}"
        depth = 1
        i = start
        while i < len(text):
            if text[i:i + len(begin_tag)] == begin_tag:
                depth += 1
                i += len(begin_tag)
            elif text[i:i + len(end_tag)] == end_tag:
                depth -= 1
                if depth == 0:
                    return i
                i += len(end_tag)
            else:
                i += 1
        raise ValueError(f"No matching \\end{{{env_name}}} found from position {start}")

    def extract_brace_arg(self, text: str, pos: int) -> tuple:
        """Extract {arg} at pos. Returns (arg_content, end_pos)."""
        pos = self._skip_ws(text, pos)
        if pos >= len(text) or text[pos] != "{":
            return None, pos
        tok = TexTokenizer(text)
        end = tok.find_matching_brace(pos)
        return text[pos + 1:end], end + 1

    def extract_bracket_arg(self, text: str, pos: int) -> tuple:
        """Extract optional [arg] at pos. Returns (arg_content, end_pos)."""
        pos = self._skip_ws(text, pos)
        if pos >= len(text) or text[pos] != "[":
            return None, pos
        tok = TexTokenizer(text)
        end = tok.find_matching_bracket(pos)
        return text[pos + 1:end], end + 1

    def _skip_ws(self, text: str, pos: int) -> int:
        while pos < len(text) and text[pos] in " \t":
            pos += 1
        return pos


# ════════════════════════════════════════════════════════════════════
# 3. BEAMER PARSER
# ════════════════════════════════════════════════════════════════════

# Standard LaTeX color names
STANDARD_COLORS = {
    "red": (220, 0, 0), "blue": (0, 0, 220), "green": (0, 128, 0),
    "black": (0, 0, 0), "white": (255, 255, 255), "gray": (128, 128, 128),
    "orange": (255, 165, 0), "purple": (128, 0, 128), "cyan": (0, 255, 255),
    "yellow": (255, 255, 0), "brown": (139, 69, 19), "magenta": (255, 0, 255),
}

# Unicode substitutions for inline math
MATH_UNICODE = {
    r"\tau": "\u03c4", r"\epsilon": "\u03b5", r"\omega": "\u03c9",
    r"\lambda": "\u03bb", r"\sigma": "\u03c3", r"\alpha": "\u03b1",
    r"\beta": "\u03b2", r"\gamma": "\u03b3", r"\delta": "\u03b4",
    r"\pi": "\u03c0", r"\mu": "\u03bc", r"\phi": "\u03c6",
    r"\sim": "\u223c", r"\pm": "\u00b1", r"\times": "\u00d7",
    r"\cdot": "\u00b7", r"\le": "\u2264", r"\leq": "\u2264",
    r"\ge": "\u2265", r"\geq": "\u2265", r"\ne": "\u2260",
    r"\neq": "\u2260", r"\approx": "\u2248",
    r"\rightarrow": "\u2192", r"\leftarrow": "\u2190",
    r"\Rightarrow": "\u21d2", r"\Leftarrow": "\u21d0",
    r"\downarrow": "\u2193", r"\uparrow": "\u2191",
    r"\blacktriangleright": "\u25b8", r"\blacktriangleleft": "\u25c2",
    r"\infty": "\u221e", r"\in": "\u2208", r"\notin": "\u2209",
    r"\subset": "\u2282", r"\supset": "\u2283",
    r"\sum": "\u2211", r"\prod": "\u220f", r"\int": "\u222b",
    r"\partial": "\u2202", r"\nabla": "\u2207",
    r"\forall": "\u2200", r"\exists": "\u2203",
    r"\mathbb{1}": "\U0001d7d9", r"\mathbbm{1}": "\U0001d7d9",
    r"\mathbf{1}": "\U0001d7cf",
}

# LaTeX font size -> approximate pt
LATEX_FONT_SIZES = {
    "tiny": 8, "scriptsize": 10, "footnotesize": 11,
    "small": 14, "normalsize": 18, "large": 22,
    "Large": 26, "LARGE": 30, "huge": 36, "Huge": 42,
}


class BeamerParser:
    """Recursive descent parser for Beamer .tex files."""

    def __init__(self, source: str):
        self.source = source
        self.tok = TexTokenizer(source)

    def parse(self) -> tuple:
        """Parse entire document. Returns (PreambleData, list[TexNode frames])."""
        # Split at \begin{document}
        doc_start = self.source.find("\\begin{document}")
        if doc_start < 0:
            raise ValueError("No \\begin{document} found")
        preamble_text = self.source[:doc_start]
        body_text = self.source[doc_start:]

        preamble = self._parse_preamble(preamble_text)
        frames = self._parse_frames(body_text)
        return preamble, frames

    # ── Preamble ──────────────────────────────────────────────────

    def _parse_preamble(self, text: str) -> PreambleData:
        p = PreambleData()

        # Aspect ratio
        m = re.search(r"\\documentclass\[([^\]]*)\]", text)
        if m:
            opts = m.group(1)
            ar = re.search(r"aspectratio=(\d+)", opts)
            if ar:
                p.aspect_ratio = ar.group(1)

        # Colors
        for m in re.finditer(r"\\definecolor\{(\w+)\}\{RGB\}\{(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\}", text):
            p.colors[m.group(1)] = (int(m.group(2)), int(m.group(3)), int(m.group(4)))

        # Beamer colors
        for m in re.finditer(r"\\setbeamercolor\{([^}]+)\}\{([^}]+)\}", text):
            elem = m.group(1)
            opts = {}
            for part in m.group(2).split(","):
                part = part.strip()
                if "=" in part:
                    k, v = part.split("=", 1)
                    opts[k.strip()] = v.strip()
            p.beamer_colors[elem] = opts

        # Title — use brace matching for nested braces
        m = re.search(r"\\title(?:\[[^\]]*\])?\{", text)
        if m:
            brace_pos = m.end() - 1
            try:
                end = TexTokenizer(text).find_matching_brace(brace_pos)
                raw_title = text[brace_pos + 1:end]
                # Preserve \\ as newline
                p.title = self._clean_title_text(raw_title)
            except ValueError:
                p.title = ""

        # Author — preserve line structure
        m = re.search(r"\\author(?:\[[^\]]*\])?\{", text)
        if m:
            brace_pos = m.end() - 1
            try:
                end = TexTokenizer(text).find_matching_brace(brace_pos)
                raw_author = text[brace_pos + 1:end]
                # Split on \\ to preserve multi-line authors
                lines = re.split(r"\\\\", raw_author)
                cleaned_lines = []
                for line in lines:
                    line = re.sub(r"\\quad\b", "    ", line)
                    line = re.sub(r"\\[a-zA-Z]+\{([^}]*)\}", r"\1", line)
                    line = re.sub(r"\\[a-zA-Z]+", " ", line)
                    line = re.sub(r"[{}]", "", line)
                    line = re.sub(r"\s+", " ", line).strip()
                    if line:
                        cleaned_lines.append(line)
                p.author = "\n".join(cleaned_lines)
            except ValueError:
                p.author = ""

        # Institute
        m = re.search(r"\\institute(?:\[[^\]]*\])?\{(.+?)\}", text, re.DOTALL)
        if m:
            p.institute = self._clean_text(m.group(1))

        # Date
        m = re.search(r"\\date\{(.*?)\}", text)
        if m:
            p.date = m.group(1).strip()

        return p

    # ── Frame splitting ───────────────────────────────────────────

    def _parse_frames(self, body: str) -> list:
        frames = []
        frame_num = 0
        i = 0
        while True:
            idx = body.find("\\begin{frame}", i)
            if idx < 0:
                break

            frame_num += 1
            # Find content start (after \begin{frame} and optional args)
            after_begin = idx + len("\\begin{frame}")
            # Check for optional [args]
            _, after_begin = self.tok.extract_bracket_arg(body, after_begin)
            # Check for {Title}
            title = None
            title_content, after_title = self.tok.extract_brace_arg(body, after_begin)
            if title_content is not None:
                title = title_content
                content_start = after_title
            else:
                content_start = after_begin

            # Find \end{frame}
            end_pos = self.tok.find_env_end(body, "frame", content_start)
            frame_body = body[content_start:end_pos].strip()

            frame_node = TexNode(
                kind="frame",
                attrs={"title": title, "slide_number": frame_num},
                children=self._parse_content(frame_body)
            )
            frames.append(frame_node)
            i = end_pos + len("\\end{frame}")

        logger.info(f"Parsed {len(frames)} frames")
        return frames

    # ── Content parser (recursive) ────────────────────────────────

    def _parse_content(self, text: str) -> list:
        """Parse frame/block/column body into a list of TexNode children."""
        nodes = []
        i = 0
        text_buf = []

        def flush_text():
            nonlocal text_buf
            t = "".join(text_buf).strip()
            if t:
                runs = self._parse_rich_text(t)
                if runs:
                    nodes.append(TexNode(kind="text", attrs={"runs": runs}))
            text_buf = []

        while i < len(text):
            # Check for environments
            env_match = re.match(r"\\begin\{(\w+)\}", text[i:])
            if env_match:
                flush_text()
                env_name = env_match.group(1)
                env_start = i + env_match.end()
                node, new_i = self._parse_env(text, env_name, i, env_start)
                if node:
                    nodes.append(node)
                i = new_i
                continue

            # Check for commands
            cmd_result = self._try_parse_command(text, i)
            if cmd_result:
                node, new_i = cmd_result
                if node:
                    flush_text()
                    nodes.append(node)
                i = new_i
                continue

            # Check for \\ (line break) — must come before \[ check
            if text[i:i+2] == "\\\\":
                # This is a LaTeX line break \\, possibly with optional [spacing]
                # Skip the \\ and any [Npt] after it
                j = i + 2
                if j < len(text) and text[j] == "[":
                    bracket_end = text.find("]", j)
                    if bracket_end >= 0:
                        j = bracket_end + 1
                text_buf.append("\n")
                i = j
                continue

            # Check for display math \[...\]
            if text[i:i+2] == "\\[":
                flush_text()
                end = text.find("\\]", i + 2)
                if end >= 0:
                    math_tex = text[i+2:end].strip()
                    nodes.append(TexNode(kind="math_display", attrs={"tex": math_tex}))
                    i = end + 2
                    continue

            # Check for comments
            if text[i] == "%" and (i == 0 or text[i-1] != "\\"):
                eol = text.find("\n", i)
                i = eol + 1 if eol >= 0 else len(text)
                continue

            text_buf.append(text[i])
            i += 1

        flush_text()
        return nodes

    def _parse_env(self, text: str, env_name: str, begin_pos: int, content_start: int) -> tuple:
        """Parse a \\begin{env}...\\end{env}. Returns (TexNode, pos_after_end)."""
        try:
            env_end = self.tok.find_env_end(text, env_name, content_start)
        except ValueError:
            logger.warning(f"Could not find \\end{{{env_name}}}")
            return None, content_start

        after_end = env_end + len(f"\\end{{{env_name}}}")

        if env_name == "columns":
            return self._build_columns(text, begin_pos, content_start, env_end), after_end
        elif env_name == "column":
            return self._build_column(text, begin_pos, content_start, env_end), after_end
        elif env_name == "block":
            return self._build_block(text, begin_pos, content_start, env_end), after_end
        elif env_name == "itemize":
            return self._build_itemize(text, content_start, env_end), after_end
        elif env_name == "enumerate":
            return self._build_enumerate(text, content_start, env_end), after_end
        elif env_name == "tikzpicture":
            body = text[content_start:env_end]
            return self._build_tikz(body), after_end
        elif env_name in ("tabular", "tabular*"):
            return self._build_tabular(text, begin_pos, content_start, env_end), after_end
        elif env_name == "beamercolorbox":
            # Skip beamercolorbox (used for narration styling)
            return None, after_end
        else:
            # Unknown env: wrap as generic
            body = text[content_start:env_end]
            return TexNode(kind="env", attrs={"name": env_name, "body": body}), after_end

    def _try_parse_command(self, text: str, pos: int) -> Optional[tuple]:
        """Try to parse a command at pos. Returns (TexNode, new_pos) or None."""
        if text[pos] != "\\":
            return None

        # \titlepage
        if text[pos:pos+10] == "\\titlepage":
            return TexNode(kind="titlepage"), pos + 10

        # \centering
        if text[pos:pos+10] == "\\centering":
            return TexNode(kind="centering"), pos + 10

        # \vspace{...}
        m = re.match(r"\\vspace\*?\{([^}]+)\}", text[pos:])
        if m:
            return TexNode(kind="vspace", attrs={"amount": m.group(1)}), pos + m.end()

        # \includegraphics[opts]{path}
        m = re.match(r"\\includegraphics\[([^\]]*)\]\{([^}]+)\}", text[pos:])
        if m:
            opts_str = m.group(1)
            path = m.group(2)
            opts = self._parse_kv_options(opts_str)
            return TexNode(kind="includegraphics", attrs={
                "path": path, "width": opts.get("width", ""),
                "height": opts.get("height", ""),
                "keepaspectratio": "keepaspectratio" in opts_str,
            }), pos + m.end()

        # \narration{...}
        if text[pos:pos+10] == "\\narration":
            brace_pos = text.find("{", pos + 10)
            if brace_pos >= 0:
                try:
                    end = TexTokenizer(text).find_matching_brace(brace_pos)
                    narr = text[brace_pos+1:end]
                    return TexNode(kind="narration", attrs={"text": narr}), end + 1
                except ValueError:
                    pass

        # \qrcode[opts]{url}
        m = re.match(r"\\qrcode(?:\[([^\]]*)\])?\{([^}]+)\}", text[pos:])
        if m:
            return TexNode(kind="qrcode", attrs={
                "height": m.group(1) or "1cm",
                "url": m.group(2),
            }), pos + m.end()

        # Skip noise commands (don't consume as text)
        for skip_cmd in [
            r"\\renewcommand\{[^}]+\}\{[^}]+\}",
            r"\\setlength\{[^}]+\}\{[^}]+\}",
            r"\\hspace\*?\{[^}]+\}",
            r"\\raggedright",
            r"\\noindent",
        ]:
            m = re.match(skip_cmd, text[pos:])
            if m:
                return None, pos + m.end()

        return None

    # ── Environment builders ──────────────────────────────────────

    def _build_columns(self, text: str, begin_pos: int, content_start: int, env_end: int) -> TexNode:
        # Extract optional [T] alignment
        after_begin = begin_pos + len("\\begin{columns}")
        align, body_start = self.tok.extract_bracket_arg(text, after_begin)
        body = text[body_start:env_end]
        children = self._parse_content(body)
        return TexNode(kind="columns", attrs={"alignment": align or "T"}, children=children)

    def _build_column(self, text: str, begin_pos: int, content_start: int, env_end: int) -> TexNode:
        after_begin = begin_pos + len("\\begin{column}")
        width_str, body_start = self.tok.extract_brace_arg(text, after_begin)
        width_frac = 0.5
        if width_str:
            m = re.search(r"([\d.]+)\\textwidth", width_str)
            if m:
                width_frac = float(m.group(1))
        body = text[body_start:env_end]
        children = self._parse_content(body)
        return TexNode(kind="column", attrs={"width_frac": width_frac}, children=children)

    def _build_block(self, text: str, begin_pos: int, content_start: int, env_end: int) -> TexNode:
        after_begin = begin_pos + len("\\begin{block}")
        title_str, body_start = self.tok.extract_brace_arg(text, after_begin)
        title = self._clean_text(title_str or "")
        body = text[body_start:env_end]
        children = self._parse_content(body)
        return TexNode(kind="block", attrs={"title": title}, children=children)

    def _build_itemize(self, text: str, content_start: int, env_end: int) -> TexNode:
        body = text[content_start:env_end].strip()
        # Remove \setlength commands at start
        body = re.sub(r"\\setlength\{[^}]+\}\{[^}]+\}", "", body).strip()
        # Detect font size switches
        font_size = None
        for sz in LATEX_FONT_SIZES:
            if body.startswith(f"\\{sz}"):
                font_size = sz
                body = body[len(f"\\{sz}"):].strip()
                break

        # Split at \item
        items = []
        parts = re.split(r"\\item(?:\[([^\]]*)\])?", body)
        # parts[0] is before first \item (usually empty), then alternating (marker, content)
        idx = 1
        while idx < len(parts):
            marker = parts[idx]  # optional [marker] content or None
            content = parts[idx + 1] if idx + 1 < len(parts) else ""
            content = content.strip()
            runs = self._parse_rich_text(content)
            items.append(TexNode(kind="item", attrs={
                "marker": marker, "runs": runs, "font_size": font_size
            }))
            idx += 2

        return TexNode(kind="itemize", attrs={"font_size": font_size}, children=items)

    def _build_enumerate(self, text: str, content_start: int, env_end: int) -> TexNode:
        body = text[content_start:env_end].strip()
        body = re.sub(r"\\setlength\{[^}]+\}\{[^}]+\}", "", body).strip()
        font_size = None
        for sz in LATEX_FONT_SIZES:
            if body.startswith(f"\\{sz}"):
                font_size = sz
                body = body[len(f"\\{sz}"):].strip()
                break
        items = []
        parts = re.split(r"\\item(?:\[([^\]]*)\])?", body)
        idx = 1
        while idx < len(parts):
            marker = parts[idx]
            content = parts[idx + 1] if idx + 1 < len(parts) else ""
            runs = self._parse_rich_text(content.strip())
            items.append(TexNode(kind="item", attrs={"marker": marker, "runs": runs}))
            idx += 2
        return TexNode(kind="enumerate", attrs={"font_size": font_size}, children=items)

    def _build_tikz(self, body: str) -> TexNode:
        """Classify and parse a tikzpicture."""
        tikz_type = self._classify_tikz(body)

        if tikz_type == "callout":
            callout_data = self._parse_tikz_callout(body)
            return TexNode(kind="tikz_callout", attrs=callout_data)
        else:
            return TexNode(kind="tikz_complex", attrs={"raw": body})

    def _classify_tikz(self, body: str) -> str:
        """Classify tikzpicture as 'callout' (single node box) or 'complex'."""
        has_foreach = "\\foreach" in body
        has_draw = bool(re.search(r"\\draw\b", body))
        has_fill = bool(re.search(r"\\fill\b", body))
        node_count = body.count("\\node")

        if has_foreach or has_fill:
            return "complex"
        if has_draw and node_count > 2:
            return "complex"
        if node_count <= 2 and not has_draw:
            return "callout"
        return "complex"

    def _parse_tikz_callout(self, body: str) -> dict:
        """Parse a simple callout tikzpicture node."""
        result = {"draw": "ubblue", "fill": "lightbg", "text_width": "0.88\\textwidth",
                  "content_runs": [TextRun(text="")]}

        # Extract node options
        m = re.search(r"\\node\[([^\]]+)\]", body)
        if m:
            opts = self._parse_kv_options(m.group(1))
            result["draw"] = opts.get("draw", "ubblue")
            result["fill"] = opts.get("fill", "lightbg")
            result["text_width"] = opts.get("text width", "0.88\\textwidth")

        # Extract text content: everything between the first { after ] and the matching }
        bracket_end = body.find("]")
        if bracket_end >= 0:
            brace_start = body.find("{", bracket_end)
            if brace_start >= 0:
                try:
                    brace_end = TexTokenizer(body).find_matching_brace(brace_start)
                    content = body[brace_start + 1:brace_end].strip()
                    result["content_runs"] = self._parse_rich_text(content)
                except ValueError:
                    pass

        return result

    def _build_tabular(self, text: str, begin_pos: int, content_start: int, env_end: int) -> TexNode:
        """Parse \\begin{tabular}{spec}...\\end{tabular}."""
        after_begin = begin_pos + len("\\begin{tabular}")
        # There may be an optional * after tabular
        if after_begin < len(text) and text[after_begin] == "*":
            after_begin += 1
        col_spec, body_start = self.tok.extract_brace_arg(text, after_begin)

        body = text[body_start:env_end].strip()

        # Split into rows by \\ (brace-aware)
        raw_rows = self._split_table_rows(body)

        rows = []
        for raw_row in raw_rows:
            raw_row = raw_row.strip()
            if not raw_row:
                continue
            # Check for rules
            if raw_row in ("\\toprule", "\\midrule", "\\bottomrule"):
                rows.append(TableRow(is_rule=True))
                continue
            # Remove trailing rule commands
            raw_row = re.sub(r"\\(toprule|midrule|bottomrule)\s*$", "", raw_row).strip()
            if not raw_row:
                continue

            # Check for \rowcolor
            rowcolor = None
            m = re.match(r"\\rowcolor\{([^}]+)\}\s*", raw_row)
            if m:
                rowcolor = m.group(1)
                raw_row = raw_row[m.end():]

            # Split cells by &
            raw_cells = raw_row.split("&")
            cells = []
            for cell_str in raw_cells:
                cell_str = cell_str.strip()
                cellcolor = None
                m2 = re.match(r"\\cellcolor\{([^}]+)\}\s*", cell_str)
                if m2:
                    cellcolor = m2.group(1)
                    cell_str = cell_str[m2.end():]

                multirow = None
                m3 = re.match(r"\\multirow\{(\d+)\}\{[^}]*\}\{", cell_str)
                if m3:
                    multirow = int(m3.group(1))
                    brace_start = cell_str.find("{", m3.start() + len(m3.group(0)) - 1)
                    if brace_start < 0:
                        brace_start = cell_str.index("{", m3.end() - 1)
                    try:
                        brace_end = TexTokenizer(cell_str).find_matching_brace(brace_start)
                        cell_str = cell_str[brace_start + 1:brace_end]
                    except ValueError:
                        cell_str = cell_str[m3.end():]

                runs = self._parse_rich_text(cell_str.strip())
                cells.append(TableCell(runs=runs, cellcolor=cellcolor, multirow=multirow))

            rows.append(TableRow(cells=cells, rowcolor=rowcolor))

        return TexNode(kind="tabular", attrs={"col_spec": col_spec, "rows": rows})

    def _split_table_rows(self, body: str) -> list:
        """Split table body by \\\\ respecting brace nesting."""
        rows = []
        depth = 0
        current = []
        i = 0
        while i < len(body):
            if body[i] == "{":
                depth += 1
            elif body[i] == "}":
                depth -= 1
            elif body[i:i+2] == "\\\\" and depth == 0:
                rows.append("".join(current).strip())
                current = []
                i += 2
                # Skip optional [spacing]
                while i < len(body) and body[i] in " \t":
                    i += 1
                if i < len(body) and body[i] == "[":
                    end = body.find("]", i)
                    if end >= 0:
                        i = end + 1
                continue
            current.append(body[i])
            i += 1
        last = "".join(current).strip()
        if last:
            rows.append(last)
        return rows

    # ── Rich text parser ──────────────────────────────────────────

    def _parse_rich_text(self, text: str) -> list:
        """Parse formatted text into list of TextRun."""
        if not text:
            return []

        # Clean up LaTeX artifacts
        text = text.replace("~", " ")
        # Handle dashes BEFORE other processing
        text = text.replace("---", "\u2014")
        text = text.replace("--", "\u2013")
        text = re.sub(r"\\,", " ", text)
        text = re.sub(r"\\!", "", text)
        text = re.sub(r"\\;", " ", text)
        text = re.sub(r"\\quad\b", "  ", text)
        text = re.sub(r"\\qquad\b", "    ", text)
        text = re.sub(r"\\hfill\b", "  ", text)
        # Handle line breaks
        text = re.sub(r"\\\\(?:\[[\d.]+\w+\])?", "\n", text)
        text = re.sub(r"\\par\b", "\n", text)
        # Remove \setlength and similar
        text = re.sub(r"\\setlength\{[^}]+\}\{[^}]+\}", "", text)
        text = re.sub(r"\\renewcommand\{[^}]+\}\{[^}]+\}", "", text)

        return self._parse_runs(text)

    def _parse_runs(self, text: str, inherit_bold=False, inherit_italic=False,
                    inherit_color=None, inherit_mono=False, inherit_size=None) -> list:
        """Recursively parse text into TextRun list."""
        runs = []
        i = 0
        buf = []

        def flush():
            nonlocal buf
            t = "".join(buf)
            if t:
                runs.append(TextRun(
                    text=t, bold=inherit_bold, italic=inherit_italic,
                    color=inherit_color, monospace=inherit_mono, size=inherit_size
                ))
            buf = []

        while i < len(text):
            ch = text[i]

            if ch == "\\" and i + 1 < len(text):
                # Try formatting commands
                cmd_match = re.match(
                    r"\\(textbf|textit|emph|textcolor|texttt|href|underline|mathbf|mathrm)\b",
                    text[i:]
                )
                if cmd_match:
                    flush()
                    cmd = cmd_match.group(1)
                    pos = i + cmd_match.end()

                    if cmd == "textcolor":
                        color_arg, pos = self.tok.extract_brace_arg(text, pos)
                        inner_arg, pos = self.tok.extract_brace_arg(text, pos)
                        if inner_arg is not None:
                            inner_runs = self._parse_runs(
                                inner_arg, inherit_bold, inherit_italic,
                                color_arg, inherit_mono, inherit_size
                            )
                            runs.extend(inner_runs)
                    elif cmd == "href":
                        url_arg, pos = self.tok.extract_brace_arg(text, pos)
                        text_arg, pos = self.tok.extract_brace_arg(text, pos)
                        if text_arg is not None:
                            inner_runs = self._parse_runs(
                                text_arg, inherit_bold, inherit_italic,
                                inherit_color, inherit_mono, inherit_size
                            )
                            for r in inner_runs:
                                r.href = url_arg
                            runs.extend(inner_runs)
                    elif cmd in ("textbf", "mathbf"):
                        arg, pos = self.tok.extract_brace_arg(text, pos)
                        if arg is not None:
                            inner_runs = self._parse_runs(
                                arg, True, inherit_italic,
                                inherit_color, inherit_mono, inherit_size
                            )
                            runs.extend(inner_runs)
                    elif cmd in ("textit", "emph"):
                        arg, pos = self.tok.extract_brace_arg(text, pos)
                        if arg is not None:
                            inner_runs = self._parse_runs(
                                arg, inherit_bold, True,
                                inherit_color, inherit_mono, inherit_size
                            )
                            runs.extend(inner_runs)
                    elif cmd == "texttt":
                        arg, pos = self.tok.extract_brace_arg(text, pos)
                        if arg is not None:
                            inner_runs = self._parse_runs(
                                arg, inherit_bold, inherit_italic,
                                inherit_color, True, inherit_size
                            )
                            runs.extend(inner_runs)
                    elif cmd == "underline":
                        arg, pos = self.tok.extract_brace_arg(text, pos)
                        if arg is not None:
                            inner_runs = self._parse_runs(
                                arg, inherit_bold, inherit_italic,
                                inherit_color, inherit_mono, inherit_size
                            )
                            runs.extend(inner_runs)
                    else:
                        arg, pos = self.tok.extract_brace_arg(text, pos)
                        if arg is not None:
                            runs.append(TextRun(text=arg, bold=inherit_bold,
                                                italic=inherit_italic, color=inherit_color))

                    i = pos
                    continue

                # Font size switches
                size_match = re.match(
                    r"\\(tiny|scriptsize|footnotesize|small|normalsize|large|Large|LARGE|huge|Huge)\b",
                    text[i:]
                )
                if size_match:
                    flush()
                    inherit_size = size_match.group(1)
                    i += size_match.end()
                    continue

                # Skip known harmless commands
                skip_match = re.match(
                    r"\\(centering|raggedright|noindent|vfill|hfill|hspace\*?\{[^}]*\}|vspace\*?\{[^}]*\})\s*",
                    text[i:]
                )
                if skip_match:
                    i += skip_match.end()
                    continue

                # \& -> &
                if text[i:i+2] == "\\&":
                    buf.append("&")
                    i += 2
                    continue

                # \--- or \-- (em-dash, en-dash)
                if text[i:i+2] == "\\-":
                    buf.append("-")
                    i += 2
                    continue

                # Other escaped chars
                if text[i+1] in "{}%$&#_~^":
                    buf.append(text[i+1])
                    i += 2
                    continue

                # Unknown command - keep as text
                cmd_m = re.match(r"\\(\w+)", text[i:])
                if cmd_m:
                    # Try extracting arg
                    cmd_name = cmd_m.group(1)
                    pos = i + cmd_m.end()
                    arg, new_pos = self.tok.extract_brace_arg(text, pos)
                    if arg is not None:
                        buf.append(arg)
                        i = new_pos
                    else:
                        buf.append(text[i:i + cmd_m.end()])
                        i += cmd_m.end()
                    continue

            elif ch == "$":
                # Inline math
                flush()
                end = text.find("$", i + 1)
                if end >= 0:
                    math_str = text[i+1:end]
                    unicode_text = self._math_to_unicode(math_str)
                    runs.append(TextRun(
                        text=unicode_text, bold=inherit_bold, italic=True,
                        color=inherit_color, size=inherit_size
                    ))
                    i = end + 1
                else:
                    buf.append("$")
                    i += 1
                continue

            elif ch == "{":
                # Scoped group - parse content recursively
                try:
                    end = TexTokenizer(text).find_matching_brace(i)
                    inner = text[i+1:end]
                    flush()
                    inner_runs = self._parse_runs(
                        inner, inherit_bold, inherit_italic,
                        inherit_color, inherit_mono, inherit_size
                    )
                    runs.extend(inner_runs)
                    i = end + 1
                    continue
                except ValueError:
                    buf.append(ch)
                    i += 1
                    continue

            buf.append(ch)
            i += 1

        flush()
        return runs

    def _math_to_unicode(self, math_str: str) -> str:
        """Convert inline math to Unicode text."""
        result = math_str
        # Sort by length (longest first) to avoid partial replacements
        for latex, uni in sorted(MATH_UNICODE.items(), key=lambda x: -len(x[0])):
            result = result.replace(latex, uni)
        # Clean remaining LaTeX
        result = re.sub(r"\\[a-zA-Z]+\{([^}]*)\}", r"\1", result)
        result = re.sub(r"\\[a-zA-Z]+", "", result)
        result = result.replace("{", "").replace("}", "")
        result = result.replace("^", "").replace("_", "")
        result = re.sub(r"\s+", " ", result).strip()
        return result

    # ── Utility ───────────────────────────────────────────────────

    def _parse_kv_options(self, opts_str: str) -> dict:
        """Parse key=value options from TikZ/includegraphics option strings."""
        opts = {}
        # Split on commas, but respect nested braces
        parts = []
        depth = 0
        current = []
        for ch in opts_str:
            if ch == "{":
                depth += 1
            elif ch == "}":
                depth -= 1
            elif ch == "," and depth == 0:
                parts.append("".join(current).strip())
                current = []
                continue
            current.append(ch)
        if current:
            parts.append("".join(current).strip())

        for part in parts:
            if "=" in part:
                k, v = part.split("=", 1)
                opts[k.strip()] = v.strip()
            elif part.strip():
                opts[part.strip()] = True
        return opts

    def _clean_text(self, text: str) -> str:
        """Remove LaTeX formatting for plain text extraction."""
        text = re.sub(r"\\textbf\{([^}]*)\}", r"\1", text)
        text = re.sub(r"\\textit\{([^}]*)\}", r"\1", text)
        text = re.sub(r"\\emph\{([^}]*)\}", r"\1", text)
        text = re.sub(r"\\texttt\{([^}]*)\}", r"\1", text)
        text = re.sub(r"\\textcolor\{[^}]*\}\{([^}]*)\}", r"\1", text)
        text = text.replace("\\\\", " ")
        text = re.sub(r"\\[a-zA-Z]+", " ", text)
        text = re.sub(r"[{}]", "", text)
        text = re.sub(r"\s+", " ", text).strip()
        return text

    def _clean_title_text(self, text: str) -> str:
        """Clean title text, preserving \\\\ as newline."""
        text = text.replace("--", "\u2013")
        text = re.sub(r"\\\\(?:\[[\d.]+\w+\])?", "\n", text)
        text = re.sub(r"\\textbf\{([^}]*)\}", r"\1", text)
        text = re.sub(r"\\textit\{([^}]*)\}", r"\1", text)
        text = re.sub(r"\\[a-zA-Z]+", " ", text)
        text = re.sub(r"[{}]", "", text)
        # Clean up whitespace per line
        lines = [re.sub(r"\s+", " ", l).strip() for l in text.split("\n")]
        return "\n".join(l for l in lines if l)


# ════════════════════════════════════════════════════════════════════
# 4. LAYOUT ENGINE
# ════════════════════════════════════════════════════════════════════

@dataclass
class LayoutBox:
    """A positioned rectangle on a slide."""
    node: TexNode
    x: float
    y: float
    width: float
    height: float
    children: list = field(default_factory=list)


class LayoutEngine:
    """Computes positions for AST nodes on slides."""

    SLIDE_W = 13.333
    SLIDE_H = 7.5
    MARGIN_L = 0.4
    MARGIN_R = 0.4
    TITLE_BAR_H = 0.65
    ACCENT_BAR_H = 0.03
    CONTENT_TOP = 0.85
    CONTENT_BOTTOM = 7.1
    BLOCK_HEADER_H = 0.45
    BLOCK_BODY_PAD = 0.15
    COL_GAP = 0.3
    ITEM_LINE_H = 0.28
    TABLE_ROW_H = 0.32
    DEFAULT_FONT_PT = 16

    def __init__(self, preamble: PreambleData):
        self.preamble = preamble
        self.content_w = self.SLIDE_W - self.MARGIN_L - self.MARGIN_R

    def layout_frame(self, frame: TexNode) -> list:
        """Layout all children of a frame. Returns list of LayoutBox."""
        is_titlepage = any(c.kind == "titlepage" for c in frame.children)
        if is_titlepage:
            return [LayoutBox(node=frame, x=0, y=0, width=self.SLIDE_W, height=self.SLIDE_H)]

        boxes, _ = self._layout_children(
            frame.children, self.MARGIN_L, self.CONTENT_TOP,
            self.content_w, self.CONTENT_BOTTOM - self.CONTENT_TOP
        )
        return boxes

    def _layout_children(self, children: list, x: float, y: float,
                         width: float, avail_h: float) -> tuple:
        """Layout children vertically. Returns (boxes, total_height)."""
        boxes = []
        cursor_y = y

        for child in children:
            if child.kind == "centering":
                continue
            if child.kind == "vspace":
                cursor_y += self._parse_dim(child.attrs.get("amount", "0pt"))
                continue
            if child.kind == "narration":
                # Narration doesn't take space
                boxes.append(LayoutBox(node=child, x=x, y=cursor_y, width=width, height=0))
                continue

            box, h = self._layout_node(child, x, cursor_y, width)
            boxes.append(box)
            cursor_y += h

        return boxes, cursor_y - y

    def _layout_node(self, node: TexNode, x: float, y: float, width: float) -> tuple:
        """Layout a single node. Returns (LayoutBox, height_used)."""
        kind = node.kind

        if kind == "columns":
            return self._layout_columns(node, x, y, width)
        elif kind == "column":
            return self._layout_column(node, x, y, width)
        elif kind == "block":
            return self._layout_block(node, x, y, width)
        elif kind == "itemize" or kind == "enumerate":
            return self._layout_itemize(node, x, y, width)
        elif kind == "tabular":
            return self._layout_tabular(node, x, y, width)
        elif kind == "tikz_callout":
            return self._layout_tikz_callout(node, x, y, width)
        elif kind == "tikz_complex":
            return self._layout_tikz_complex(node, x, y, width)
        elif kind == "includegraphics":
            return self._layout_image(node, x, y, width)
        elif kind == "math_display":
            h = 0.7
            return LayoutBox(node=node, x=x, y=y, width=width, height=h), h
        elif kind == "qrcode":
            h = 1.0
            return LayoutBox(node=node, x=x, y=y, width=width, height=h), h
        elif kind == "text":
            return self._layout_text(node, x, y, width)
        else:
            return LayoutBox(node=node, x=x, y=y, width=width, height=0.1), 0.1

    def _layout_columns(self, node: TexNode, x: float, y: float, width: float) -> tuple:
        columns = [c for c in node.children if c.kind == "column"]
        if not columns:
            return LayoutBox(node=node, x=x, y=y, width=width, height=0.1), 0.1

        total_frac = sum(c.attrs.get("width_frac", 0.5) for c in columns)
        gap = self.COL_GAP
        total_gap = gap * (len(columns) - 1)
        usable_w = width - total_gap

        col_boxes = []
        col_x = x
        max_h = 0

        for col in columns:
            frac = col.attrs.get("width_frac", 0.5)
            col_w = usable_w * (frac / total_frac)
            child_boxes, col_h = self._layout_children(
                col.children, col_x, y, col_w, self.SLIDE_H - y
            )
            cb = LayoutBox(node=col, x=col_x, y=y, width=col_w, height=col_h,
                           children=child_boxes)
            col_boxes.append(cb)
            col_x += col_w + gap
            max_h = max(max_h, col_h)

        return LayoutBox(node=node, x=x, y=y, width=width, height=max_h,
                         children=col_boxes), max_h

    def _layout_column(self, node: TexNode, x: float, y: float, width: float) -> tuple:
        child_boxes, h = self._layout_children(node.children, x, y, width, self.SLIDE_H - y)
        return LayoutBox(node=node, x=x, y=y, width=width, height=h, children=child_boxes), h

    def _layout_block(self, node: TexNode, x: float, y: float, width: float) -> tuple:
        header_h = self.BLOCK_HEADER_H
        pad = self.BLOCK_BODY_PAD
        body_y = y + header_h + pad
        body_w = width - 2 * pad

        child_boxes, content_h = self._layout_children(
            node.children, x + pad, body_y, body_w, self.SLIDE_H - body_y
        )
        total_h = header_h + content_h + 2 * pad
        box = LayoutBox(node=node, x=x, y=y, width=width, height=total_h,
                        children=child_boxes)
        return box, total_h

    def _layout_itemize(self, node: TexNode, x: float, y: float, width: float) -> tuple:
        n_items = len(node.children)
        font_size = node.attrs.get("font_size")
        line_h = self.ITEM_LINE_H
        if font_size in ("scriptsize", "footnotesize", "small"):
            line_h = 0.24

        total_h = 0
        for item in node.children:
            runs = item.attrs.get("runs", [])
            text_len = sum(len(r.text) for r in runs)
            chars_per_line = max(1, int(width * 72 / self.DEFAULT_FONT_PT * 1.5))
            n_lines = max(1, (text_len // chars_per_line) + 1)
            total_h += n_lines * line_h

        box = LayoutBox(node=node, x=x, y=y, width=width, height=total_h)
        return box, total_h

    def _layout_tabular(self, node: TexNode, x: float, y: float, width: float) -> tuple:
        rows = node.attrs.get("rows", [])
        data_rows = [r for r in rows if not r.is_rule]
        h = len(data_rows) * self.TABLE_ROW_H
        return LayoutBox(node=node, x=x, y=y, width=width, height=h), h

    def _layout_tikz_callout(self, node: TexNode, x: float, y: float, width: float) -> tuple:
        runs = node.attrs.get("content_runs", [])
        text_len = sum(len(r.text) for r in runs)
        lines = max(1, text_len // 80 + 1)
        h = max(0.5, lines * 0.3 + 0.2)
        return LayoutBox(node=node, x=x, y=y, width=width, height=h), h

    def _layout_tikz_complex(self, node: TexNode, x: float, y: float, width: float) -> tuple:
        h = 2.5
        return LayoutBox(node=node, x=x, y=y, width=width, height=h), h

    def _layout_image(self, node: TexNode, x: float, y: float, width: float) -> tuple:
        h = 3.5  # default
        height_str = node.attrs.get("height", "")
        if height_str:
            h = self._parse_dim(height_str)
        width_str = node.attrs.get("width", "")
        if width_str:
            img_w = self._parse_dim_relative(width_str, width)
            h = min(h, img_w * 0.6)  # rough aspect ratio
        return LayoutBox(node=node, x=x, y=y, width=width, height=h), h

    def _layout_text(self, node: TexNode, x: float, y: float, width: float) -> tuple:
        runs = node.attrs.get("runs", [])
        text_len = sum(len(r.text) for r in runs)
        font_pt = self.DEFAULT_FONT_PT
        # Check if runs have a size override
        for r in runs:
            if r.size:
                font_pt = LATEX_FONT_SIZES.get(r.size, font_pt)
                break
        chars_per_line = max(1, int(width * 72 / font_pt * 1.6))
        n_lines = max(1, (text_len // chars_per_line) + 1)
        h = n_lines * (font_pt / 72 * 1.4)
        return LayoutBox(node=node, x=x, y=y, width=width, height=h), h

    def _parse_dim(self, dim_str: str) -> float:
        """Convert LaTeX dimension to inches."""
        dim_str = dim_str.strip()
        if not dim_str:
            return 0
        try:
            if dim_str.endswith("pt"):
                return float(dim_str[:-2]) / 72.0
            elif dim_str.endswith("cm"):
                return float(dim_str[:-2]) / 2.54
            elif dim_str.endswith("mm"):
                return float(dim_str[:-2]) / 25.4
            elif dim_str.endswith("in"):
                return float(dim_str[:-2])
            elif dim_str.endswith("em"):
                return float(dim_str[:-2]) * 0.16
            else:
                return float(dim_str) / 72.0
        except ValueError:
            return 0.1

    def _parse_dim_relative(self, dim_str: str, parent_width: float) -> float:
        """Parse dimension that may be relative to textwidth."""
        m = re.match(r"([\d.]+)\\textwidth", dim_str)
        if m:
            return float(m.group(1)) * parent_width
        return self._parse_dim(dim_str)


# ════════════════════════════════════════════════════════════════════
# 5. PPTX RENDERER
# ════════════════════════════════════════════════════════════════════

class PptxRenderer:
    """Converts LayoutBox trees into python-pptx slides."""

    def __init__(self, preamble: PreambleData, tex_dir: Path,
                 img_dir: Path, pdf_path: Path, cache_dir: Path):
        self.preamble = preamble
        self.tex_dir = tex_dir
        self.img_dir = img_dir
        self.pdf_path = pdf_path
        self.cache_dir = cache_dir
        self.total_slides = 0

        # Build color lookup
        self.colors = dict(STANDARD_COLORS)
        self.colors.update(preamble.colors)

    # ── Color resolution ──────────────────────────────────────────

    def resolve_color(self, spec: str) -> RGBColor:
        """Resolve color spec like 'ubblue', 'ubblue!10', 'red!70!black'."""
        if not spec:
            return RGBColor(0, 0, 0)

        spec = spec.strip()

        if "!" not in spec:
            rgb = self.colors.get(spec, (0, 0, 0))
            return RGBColor(*rgb)

        parts = spec.split("!")
        base_name = parts[0]
        base = self.colors.get(base_name, (0, 0, 0))

        if len(parts) >= 2:
            try:
                pct = int(parts[1]) / 100.0
            except ValueError:
                return RGBColor(*base)

            mix_name = parts[2] if len(parts) > 2 else "white"
            mix = self.colors.get(mix_name, (255, 255, 255))

            r = int(base[0] * pct + mix[0] * (1 - pct))
            g = int(base[1] * pct + mix[1] * (1 - pct))
            b = int(base[2] * pct + mix[2] * (1 - pct))
            return RGBColor(min(255, r), min(255, g), min(255, b))

        return RGBColor(*base)

    # ── External rendering helpers ────────────────────────────────

    def _ensure_cache(self):
        self.cache_dir.mkdir(exist_ok=True)

    def pdf_to_png(self, pdf_path: Path, dpi: int = 300) -> Path:
        """Convert a PDF figure to PNG."""
        self._ensure_cache()
        out = self.cache_dir / (pdf_path.stem + ".png")
        if out.exists():
            return out
        try:
            subprocess.run(
                ["sips", "-s", "format", "png", "-s", "dpiWidth", str(dpi),
                 "-s", "dpiHeight", str(dpi), str(pdf_path), "--out", str(out)],
                check=True, capture_output=True
            )
            if out.exists():
                return out
        except (subprocess.CalledProcessError, FileNotFoundError):
            pass
        try:
            pfx = str(self.cache_dir / pdf_path.stem)
            subprocess.run(
                ["pdftoppm", "-png", "-r", str(dpi), "-singlefile", str(pdf_path), pfx],
                check=True, capture_output=True
            )
            expected = Path(pfx + ".png")
            if expected.exists():
                expected.rename(out)
                return out
        except (subprocess.CalledProcessError, FileNotFoundError):
            pass
        raise RuntimeError(f"Cannot convert {pdf_path} to PNG")

    def extract_pdf_page(self, page_num: int, name: str, dpi: int = 300) -> Optional[Path]:
        """Extract a single page from the presentation PDF."""
        if not self.pdf_path.exists():
            return None
        self._ensure_cache()
        out = self.cache_dir / f"{name}.png"
        if out.exists():
            return out
        try:
            pfx = str(self.cache_dir / name)
            subprocess.run(
                ["pdftoppm", "-png", "-r", str(dpi), "-singlefile",
                 "-f", str(page_num), "-l", str(page_num),
                 str(self.pdf_path), pfx],
                check=True, capture_output=True
            )
            expected = Path(pfx + ".png")
            if expected.exists():
                expected.rename(out)
                return out
        except (subprocess.CalledProcessError, FileNotFoundError):
            pass
        return None

    def render_math_png(self, tex: str, fontsize: int = 18, name: str = "eq") -> Path:
        """Render LaTeX math to PNG via matplotlib."""
        self._ensure_cache()
        safe = re.sub(r"[^a-zA-Z0-9]", "_", name + tex[:20])[:50]
        out = self.cache_dir / f"{safe}.png"
        if out.exists():
            return out
        # Clean tex for matplotlib compatibility
        clean_tex = tex
        # Replace commands matplotlib doesn't know
        clean_tex = clean_tex.replace(r"\mathbbm", r"\mathbb")
        clean_tex = clean_tex.replace(r"\le ", r"\leq ")
        clean_tex = clean_tex.replace(r"\ge ", r"\geq ")
        # Remove \big, \Big etc.
        clean_tex = re.sub(r"\\[Bb]ig[lrm]?(?=[\(\)\[\]\{\}|])", "", clean_tex)
        try:
            fig, ax = plt.subplots(figsize=(10, 1.2))
            ax.axis("off")
            ax.text(0.5, 0.5, f"${clean_tex}$", fontsize=fontsize,
                    ha="center", va="center", transform=ax.transAxes)
            fig.savefig(str(out), dpi=200, bbox_inches="tight",
                        pad_inches=0.03, transparent=True)
            plt.close(fig)
        except Exception:
            plt.close("all")
            # Fallback: render as plain text
            fig, ax = plt.subplots(figsize=(10, 1.2))
            ax.axis("off")
            plain = re.sub(r"\\[a-zA-Z]+", " ", tex)
            plain = plain.replace("{", "").replace("}", "")
            ax.text(0.5, 0.5, plain, fontsize=fontsize,
                    ha="center", va="center", transform=ax.transAxes)
            fig.savefig(str(out), dpi=200, bbox_inches="tight",
                        pad_inches=0.03, transparent=True)
            plt.close(fig)
        return out

    def render_qr_png(self, url: str) -> Path:
        """Generate a QR code PNG."""
        self._ensure_cache()
        out = self.cache_dir / "qr_pres.png"
        if out.exists():
            return out
        try:
            import qrcode
            qr = qrcode.make(url)
            qr.save(str(out))
        except ImportError:
            fig, ax = plt.subplots(figsize=(2, 2))
            ax.text(0.5, 0.5, "QR", fontsize=20, ha="center", va="center")
            ax.axis("off")
            fig.savefig(str(out), dpi=150, bbox_inches="tight")
            plt.close(fig)
        return out

    def _get_aspect(self, path: Path) -> float:
        with Image.open(path) as im:
            return im.width / im.height

    def _resolve_img(self, path: Path) -> Path:
        """If PDF, return cached PNG."""
        if str(path).lower().endswith(".pdf"):
            png = self.cache_dir / (path.stem + ".png")
            if png.exists():
                return png
        return path

    # ── PPTX shape helpers ────────────────────────────────────────

    def _new_slide(self, prs):
        return prs.slides.add_slide(prs.slide_layouts[6])

    def _add_rect(self, slide, l, t, w, h, color, line_color=None, line_w=None):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(l), Inches(t), Inches(w), Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        if line_color:
            s.line.color.rgb = line_color
            if line_w:
                s.line.width = Pt(line_w)
        else:
            s.line.fill.background()
        return s

    def _add_rrect(self, slide, l, t, w, h, fill, line_color=None, line_w=1):
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(l), Inches(t), Inches(w), Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        if line_color:
            s.line.color.rgb = line_color
            s.line.width = Pt(line_w)
        else:
            s.line.fill.background()
        try:
            s.adjustments[0] = 0.08
        except (IndexError, KeyError):
            pass
        return s

    def _add_textbox(self, slide, l, t, w, h, wrap=True):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tb.text_frame.word_wrap = wrap
        return tb.text_frame

    def _apply_runs(self, tf, runs: list, default_size=16, default_color=None, alignment=None):
        """Apply a list of TextRun to a text frame."""
        if alignment:
            tf.paragraphs[0].alignment = alignment

        if not runs:
            return

        p = tf.paragraphs[0]
        first = True
        for run in runs:
            text = run.text
            if not text:
                continue

            # Handle newlines by creating new paragraphs
            lines = text.split("\n")
            for li, line in enumerate(lines):
                if li > 0:
                    p = tf.add_paragraph()
                    if alignment:
                        p.alignment = alignment
                    first = True

                if not line and li > 0:
                    continue

                if not first and li == 0:
                    pass  # continue on same paragraph

                r = p.add_run()
                r.text = line
                r.font.size = Pt(LATEX_FONT_SIZES.get(run.size, default_size)
                                 if run.size else default_size)
                r.font.bold = run.bold
                r.font.italic = run.italic
                if run.color:
                    r.font.color.rgb = self.resolve_color(run.color)
                elif default_color:
                    r.font.color.rgb = default_color
                if run.monospace:
                    r.font.name = "Consolas"
                if run.href:
                    r.hyperlink.address = run.href
                first = False

    def _simple_text(self, slide, l, t, w, h, text, size=18, color=None, bold=False,
                     align=PP_ALIGN.LEFT, font=None):
        tf = self._add_textbox(slide, l, t, w, h)
        tf.paragraphs[0].alignment = align
        r = tf.paragraphs[0].add_run()
        r.text = text
        r.font.size = Pt(size)
        if color:
            r.font.color.rgb = color
        r.font.bold = bold
        if font:
            r.font.name = font
        return tf

    # ── Main rendering ────────────────────────────────────────────

    def render_presentation(self, frames: list, all_layouts: list) -> Presentation:
        """Build the full PPTX."""
        prs = Presentation()
        prs.slide_width = Inches(LayoutEngine.SLIDE_W)
        prs.slide_height = Inches(LayoutEngine.SLIDE_H)
        self.total_slides = len(frames)

        ubblue = self.resolve_color("ubblue") if "ubblue" in self.colors else RGBColor(0, 87, 160)
        accent = self.resolve_color("accent") if "accent" in self.colors else RGBColor(0, 150, 136)

        for frame, boxes in zip(frames, all_layouts):
            slide = self._new_slide(prs)
            slide_num = frame.attrs.get("slide_number", 0)
            title = frame.attrs.get("title")

            # Check if title page
            is_titlepage = any(c.kind == "titlepage" for c in frame.children)

            if is_titlepage:
                self._render_title_page(slide, frame)
            else:
                # Frame title bar
                if title:
                    self._add_rect(slide, 0, 0, LayoutEngine.SLIDE_W, 0.65, ubblue)
                    self._add_rect(slide, 0, 0.65, LayoutEngine.SLIDE_W, 0.03, accent)
                    self._simple_text(slide, 0.4, 0.1, 12.5, 0.5,
                                      self._clean_title(title),
                                      size=28, color=RGBColor(255, 255, 255), bold=True)

                # Render content
                self._render_boxes(slide, boxes, frame)

            # Slide number
            ubgray = self.resolve_color("ubgray") if "ubgray" in self.colors else RGBColor(85, 85, 85)
            self._simple_text(slide, 12.9, 7.1, 0.8, 0.3,
                              f"{slide_num}/{self.total_slides}",
                              size=11, color=ubgray, align=PP_ALIGN.RIGHT)

        return prs

    def _clean_title(self, title: str) -> str:
        """Clean LaTeX from frame title for display."""
        # Handle dashes
        title = title.replace("---", "\u2014")
        title = title.replace("--", "\u2013")
        # Handle math
        title = re.sub(r"\$([^$]*)\$", lambda m: self._math_to_unicode_simple(m.group(1)), title)
        # Handle overline
        title = re.sub(r"\\overline\{\\text\{([^}]*)\}\}", lambda m: m.group(1) + "\u0305", title)
        # Handle \text{} inside and outside math
        title = re.sub(r"\\text\{([^}]*)\}", r"\1", title)
        # Strip formatting commands
        title = re.sub(r"\\small\b\s*", "", title)
        title = re.sub(r"\\textbf\{([^}]*)\}", r"\1", title)
        title = re.sub(r"\\textit\{([^}]*)\}", r"\1", title)
        title = re.sub(r"\\emph\{([^}]*)\}", r"\1", title)
        title = re.sub(r"\\[a-zA-Z]+\{([^}]*)\}", r"\1", title)
        title = re.sub(r"[{}]", "", title)
        title = re.sub(r"\\[a-zA-Z]+", "", title)
        title = title.replace("~", " ")
        title = title.replace("\\", "")
        title = re.sub(r"\s+", " ", title).strip()
        # Fix common Unicode
        title = title.replace("@", "@")
        return title

    def _math_to_unicode_simple(self, s: str) -> str:
        for k, v in MATH_UNICODE.items():
            s = s.replace(k, v)
        s = re.sub(r"\\[a-zA-Z]+", "", s)
        s = s.replace("{", "").replace("}", "")
        return s

    def _render_title_page(self, slide, frame: TexNode):
        """Render the title page slide matching Madrid beamer theme."""
        ubblue = self.resolve_color("ubblue") if "ubblue" in self.colors else RGBColor(0, 87, 160)
        WHITE = RGBColor(255, 255, 255)

        # Full blue background
        self._add_rect(slide, 0, 0, LayoutEngine.SLIDE_W, LayoutEngine.SLIDE_H, ubblue)

        # Title — preserve line breaks from preamble
        tf = self._add_textbox(slide, 1.0, 1.3, 11.3, 1.5)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_lines = self.preamble.title.split("\n")
        for i, line in enumerate(title_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = line
            r.font.size = Pt(36)
            r.font.color.rgb = WHITE
            r.font.bold = True

        # Authors — split on \n for multi-line
        author_lines = self.preamble.author.split("\n")
        tf = self._add_textbox(slide, 1.0, 3.2, 11.3, 0.7)
        for i, line in enumerate(author_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = line
            r.font.size = Pt(14)
            r.font.color.rgb = RGBColor(220, 225, 235)

        # Institute
        self._simple_text(slide, 1.0, 4.1, 11.3, 0.4,
                          self.preamble.institute,
                          size=13, color=RGBColor(180, 195, 220), align=PP_ALIGN.CENTER)

        # Process frame children for subtitle, github link, QR code
        subtitle_placed = False
        github_placed = False
        qr_placed = False

        for child in frame.children:
            if child.kind == "text":
                runs = child.attrs.get("runs", [])

                # Scan runs for subtitle text and github link separately
                for run in runs:
                    run_text = run.text.strip()
                    if not run_text:
                        continue

                    # Check if this run is a github link (has href or contains github.com)
                    if not github_placed and (
                        (run.href and "github" in run.href.lower())
                        or "github.com" in run_text.lower()
                    ):
                        link_text = run_text
                        self._simple_text(slide, 6.5, 5.9, 3.5, 0.3,
                                          link_text, size=12,
                                          color=WHITE, font="Consolas")
                        github_placed = True
                        continue

                    # Check if this run is the subtitle
                    if not subtitle_placed and (
                        "differentiable" in run_text.lower()
                        or "drop-in" in run_text.lower()
                        or "replacement" in run_text.lower()
                    ):
                        # Clean up: remove trailing newlines and github text
                        subtitle = run_text.split("\n")[0].strip()
                        if subtitle:
                            self._simple_text(slide, 1.0, 5.0, 11.3, 0.4,
                                              subtitle, size=14,
                                              color=RGBColor(180, 230, 220),
                                              align=PP_ALIGN.CENTER)
                            subtitle_placed = True

            elif child.kind == "qrcode" and not qr_placed:
                qr_path = self.render_qr_png(child.attrs["url"])
                slide.shapes.add_picture(str(qr_path),
                                         Inches(5.3), Inches(5.65),
                                         Inches(0.9), Inches(0.9))
                qr_placed = True

    # ── Box rendering ─────────────────────────────────────────────

    def _render_boxes(self, slide, boxes: list, frame: TexNode):
        """Render all layout boxes on a slide."""
        for box in boxes:
            self._render_box(slide, box, frame)

    def _render_box(self, slide, box: LayoutBox, frame: TexNode):
        """Render a single layout box."""
        node = box.node
        kind = node.kind

        if kind == "columns":
            for child_box in box.children:
                self._render_box(slide, child_box, frame)
        elif kind == "column":
            for child_box in box.children:
                self._render_box(slide, child_box, frame)
        elif kind == "block":
            self._render_block(slide, box)
        elif kind == "itemize" or kind == "enumerate":
            self._render_itemize(slide, box)
        elif kind == "tabular":
            self._render_tabular(slide, box)
        elif kind == "tikz_callout":
            self._render_tikz_callout(slide, box)
        elif kind == "tikz_complex":
            self._render_tikz_complex(slide, box, frame)
        elif kind == "includegraphics":
            self._render_image(slide, box)
        elif kind == "math_display":
            self._render_math_display(slide, box)
        elif kind == "qrcode":
            self._render_qrcode(slide, box)
        elif kind == "text":
            self._render_text(slide, box)
        elif kind == "narration":
            self._render_narration(slide, node)

    def _render_block(self, slide, box: LayoutBox):
        """Render a beamer block: blue header + light bg body."""
        node = box.node
        title = node.attrs.get("title", "")
        x, y, w, h = box.x, box.y, box.width, box.height
        header_h = LayoutEngine.BLOCK_HEADER_H

        ubblue = self.resolve_color("ubblue") if "ubblue" in self.colors else RGBColor(0, 87, 160)
        lightbg = self.resolve_color("lightbg") if "lightbg" in self.colors else RGBColor(245, 247, 250)

        # Body background
        self._add_rrect(slide, x, y, w, h, lightbg, RGBColor(210, 210, 210), 1)
        # Header bar
        self._add_rect(slide, x, y, w, header_h, ubblue)
        # Header text
        self._simple_text(slide, x + 0.18, y + 0.05, w - 0.36, header_h,
                          title, size=20, color=RGBColor(255, 255, 255))

        # Render children
        for child_box in box.children:
            self._render_box(slide, child_box, TexNode(kind="frame"))

    def _render_itemize(self, slide, box: LayoutBox):
        """Render an itemize or enumerate list."""
        node = box.node
        x, y, w, h = box.x, box.y, box.width, max(box.height, 0.3)

        font_size_name = node.attrs.get("font_size")
        font_size = LATEX_FONT_SIZES.get(font_size_name, 16) if font_size_name else 16
        is_enum = node.kind == "enumerate"

        tf = self._add_textbox(slide, x, y, w, h)

        ubgray = self.resolve_color("ubgray") if "ubgray" in self.colors else RGBColor(85, 85, 85)

        for i, item in enumerate(node.children):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

            # Bullet/number marker
            marker = item.attrs.get("marker")
            if is_enum:
                marker_text = f"{i+1}. "
            elif marker:
                # Custom marker from \item[marker]
                marker_text = self._extract_marker_text(marker) + " "
            else:
                marker_text = "\u2022 "

            r = p.add_run()
            r.text = marker_text
            r.font.size = Pt(font_size)
            r.font.color.rgb = ubgray

            # Item content
            runs = item.attrs.get("runs", [])
            for run in runs:
                if not run.text:
                    continue
                r = p.add_run()
                r.text = run.text
                r.font.size = Pt(LATEX_FONT_SIZES.get(run.size, font_size) if run.size else font_size)
                r.font.bold = run.bold
                r.font.italic = run.italic
                if run.color:
                    r.font.color.rgb = self.resolve_color(run.color)
                else:
                    r.font.color.rgb = ubgray
                if run.monospace:
                    r.font.name = "Consolas"

    def _extract_marker_text(self, marker: str) -> str:
        """Extract display text from a custom item marker like \\textcolor{orange}{$\\blacktriangleright$}."""
        # Try to get unicode from math
        for k, v in MATH_UNICODE.items():
            marker = marker.replace(k, v)
        # Strip LaTeX commands
        marker = re.sub(r"\\textcolor\{[^}]*\}\{([^}]*)\}", r"\1", marker)
        marker = re.sub(r"\\[a-zA-Z]+\{([^}]*)\}", r"\1", marker)
        marker = re.sub(r"\\[a-zA-Z]+", "", marker)
        marker = marker.replace("{", "").replace("}", "").replace("$", "")
        return marker.strip() or "\u25b8"

    def _render_tabular(self, slide, box: LayoutBox):
        """Render a table."""
        node = box.node
        rows = node.attrs.get("rows", [])
        data_rows = [r for r in rows if not r.is_rule]
        if not data_rows:
            return

        x, y, w = box.x, box.y, box.width
        nr = len(data_rows)
        nc = max(len(r.cells) for r in data_rows)
        rh = LayoutEngine.TABLE_ROW_H

        ubblue = self.resolve_color("ubblue") if "ubblue" in self.colors else RGBColor(0, 87, 160)
        hlrow = self.resolve_color("ubblue!10") if "ubblue" in self.colors else RGBColor(219, 228, 245)

        ts = slide.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(nr * rh))
        tbl = ts.table

        # Determine if first row is a header (has bold formatting)
        first_row = data_rows[0]
        is_header = any(any(r.bold for r in c.runs) for c in first_row.cells if c.runs)

        for i, row in enumerate(data_rows):
            is_first = (i == 0 and is_header)

            for j, cell in enumerate(row.cells):
                if j >= nc:
                    break
                c = tbl.cell(i, j)

                # Cell text
                cell_text = " ".join(r.text for r in cell.runs).strip()
                c.text = cell_text

                # Styling
                if is_first:
                    c.fill.solid()
                    c.fill.fore_color.rgb = ubblue
                elif row.rowcolor:
                    c.fill.solid()
                    c.fill.fore_color.rgb = self.resolve_color(row.rowcolor)
                elif cell.cellcolor:
                    c.fill.solid()
                    c.fill.fore_color.rgb = self.resolve_color(cell.cellcolor)

                for p in c.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                    for r in p.runs:
                        r.font.size = Pt(11)
                        if is_first:
                            r.font.bold = True
                            r.font.color.rgb = RGBColor(255, 255, 255)
                            r.font.size = Pt(12)
                        else:
                            has_bold = any(run.bold for run in cell.runs)
                            r.font.bold = has_bold or bool(row.rowcolor)
                            r.font.color.rgb = RGBColor(0, 0, 0)

    def _render_tikz_callout(self, slide, box: LayoutBox):
        """Render a simple TikZ callout as a native PPTX rounded rectangle."""
        node = box.node
        draw_color = node.attrs.get("draw", "ubblue")
        fill_color = node.attrs.get("fill", "lightbg")
        content_runs = node.attrs.get("content_runs", [])

        x, y, w, h = box.x, box.y, box.width, max(box.height, 0.5)

        # Indent callout slightly
        margin = max(0, (w - min(w, 11.0)) / 2)
        cx = x + margin + 0.3
        cw = w - 2 * margin - 0.6

        border_rgb = self.resolve_color(draw_color)
        fill_rgb = self.resolve_color(fill_color)

        self._add_rrect(slide, cx, y, cw, h, fill_rgb, border_rgb, 2)

        tf = self._add_textbox(slide, cx + 0.2, y + 0.05, cw - 0.4, h - 0.1)
        self._apply_runs(tf, content_runs, default_size=15, alignment=PP_ALIGN.CENTER)

    def _render_tikz_complex(self, slide, box: LayoutBox, frame: TexNode):
        """Render complex TikZ by extracting from compiled PDF."""
        slide_num = frame.attrs.get("slide_number", 1)
        name = f"tikz_slide{slide_num}"
        png = self.extract_pdf_page(slide_num, name)
        if png and png.exists():
            # Embed the full slide page, cropped to the tikz area
            x, y, w, h = box.x, box.y, box.width, box.height
            try:
                aspect = self._get_aspect(png)
                img_w = min(w, h * aspect)
                img_h = img_w / aspect
                img_x = x + (w - img_w) / 2
                slide.shapes.add_picture(str(png), Inches(img_x), Inches(y),
                                         Inches(img_w), Inches(img_h))
            except Exception as e:
                logger.warning(f"Could not render complex TikZ on slide {slide_num}: {e}")
        else:
            # Fallback: show placeholder
            self._simple_text(slide, box.x, box.y, box.width, box.height,
                              "[Complex TikZ diagram - compile PDF for rendering]",
                              size=12, color=RGBColor(150, 150, 150), align=PP_ALIGN.CENTER)

    def _render_image(self, slide, box: LayoutBox):
        """Render an \\includegraphics image."""
        node = box.node
        img_path_str = node.attrs.get("path", "")
        x, y, w, h = box.x, box.y, box.width, box.height

        # Resolve path relative to tex directory
        img_path = (self.tex_dir / img_path_str).resolve()
        if not img_path.exists():
            img_path = (self.img_dir / Path(img_path_str).name).resolve()

        # Convert PDF to PNG if needed
        if img_path.suffix.lower() == ".pdf":
            try:
                img_path = self.pdf_to_png(img_path)
            except RuntimeError:
                logger.warning(f"Could not convert {img_path} to PNG")
                return

        img_path = self._resolve_img(img_path)
        if not img_path.exists():
            logger.warning(f"Image not found: {img_path}")
            return

        try:
            aspect = self._get_aspect(img_path)
            # Parse width/height constraints
            max_w = w
            max_h = h
            width_str = node.attrs.get("width", "")
            height_str = node.attrs.get("height", "")

            if width_str:
                m = re.match(r"([\d.]+)\\textwidth", width_str)
                if m:
                    max_w = float(m.group(1)) * w

            if height_str:
                dim = height_str.replace("cm", "")
                try:
                    max_h = float(dim) / 2.54
                except ValueError:
                    pass

            img_w = max_w
            img_h = img_w / aspect
            if img_h > max_h:
                img_h = max_h
                img_w = img_h * aspect

            img_x = x + (w - img_w) / 2  # center
            slide.shapes.add_picture(str(img_path), Inches(img_x), Inches(y),
                                     Inches(img_w), Inches(img_h))
        except Exception as e:
            logger.warning(f"Error adding image {img_path}: {e}")

    def _render_math_display(self, slide, box: LayoutBox):
        """Render display math as a PNG image."""
        node = box.node
        tex = node.attrs.get("tex", "")
        x, y, w, h = box.x, box.y, box.width, box.height

        try:
            eq_path = self.render_math_png(tex, fontsize=18, name="display")
            if eq_path.exists():
                aspect = self._get_aspect(eq_path)
                img_w = min(w * 0.9, h * aspect)
                img_h = img_w / aspect
                img_x = x + (w - img_w) / 2
                slide.shapes.add_picture(str(eq_path), Inches(img_x), Inches(y),
                                         Inches(img_w), Inches(img_h))
        except Exception as e:
            logger.warning(f"Math rendering failed: {e}")
            # Fallback: show as text
            self._simple_text(slide, x, y, w, h, f"[{tex}]",
                              size=14, color=RGBColor(0, 0, 0), align=PP_ALIGN.CENTER)

    def _render_qrcode(self, slide, box: LayoutBox):
        """Render a QR code."""
        node = box.node
        url = node.attrs.get("url", "")
        x, y, w, h = box.x, box.y, box.width, box.height

        qr_path = self.render_qr_png(url)
        qr_size = min(1.0, h)
        slide.shapes.add_picture(str(qr_path), Inches(x), Inches(y),
                                 Inches(qr_size), Inches(qr_size))

    def _render_text(self, slide, box: LayoutBox):
        """Render a text node."""
        node = box.node
        runs = node.attrs.get("runs", [])
        x, y, w, h = box.x, box.y, box.width, max(box.height, 0.3)

        if not runs or not any(r.text.strip() for r in runs):
            return

        ubgray = self.resolve_color("ubgray") if "ubgray" in self.colors else RGBColor(85, 85, 85)
        tf = self._add_textbox(slide, x, y, w, h)
        self._apply_runs(tf, runs, default_size=17, default_color=ubgray)

    def _render_narration(self, slide, node: TexNode):
        """Add narration text as speaker notes."""
        text = node.attrs.get("text", "")
        if not text:
            return
        # Clean LaTeX from narration
        text = re.sub(r"\\[a-zA-Z]+\{([^}]*)\}", r"\1", text)
        text = text.replace("---", "\u2014")
        text = text.replace("--", "\u2013")
        text = re.sub(r"\\[a-zA-Z]+", "", text)
        text = text.replace("{", "").replace("}", "")
        text = re.sub(r"\s+", " ", text).strip()

        try:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.text = text
        except Exception:
            pass


# ════════════════════════════════════════════════════════════════════
# 6. MAIN
# ════════════════════════════════════════════════════════════════════

def convert_pdf_figures(img_dir: Path, cache_dir: Path):
    """Pre-convert all PDF figures to PNG."""
    cache_dir.mkdir(exist_ok=True)
    renderer = PptxRenderer(PreambleData(), img_dir.parent, img_dir, Path("dummy"), cache_dir)

    pdf_files = list(img_dir.glob("*.pdf"))
    for pdf in pdf_files:
        try:
            png = renderer.pdf_to_png(pdf)
            print(f"  + {pdf.name} -> {png.name}")
        except RuntimeError as e:
            print(f"  x {pdf.name}: {e}")


def main():
    parser = argparse.ArgumentParser(
        description="Convert a Beamer LaTeX presentation to PPTX"
    )
    parser.add_argument("input", nargs="?", default="presentation.tex",
                        help="Input .tex file (default: presentation.tex)")
    parser.add_argument("-o", "--output", default=None,
                        help="Output .pptx file (default: <input>.pptx)")
    parser.add_argument("--pdf", default=None,
                        help="Compiled PDF for complex TikZ fallback")
    parser.add_argument("--dpi", type=int, default=300,
                        help="DPI for image rendering (default: 300)")
    parser.add_argument("-v", "--verbose", action="store_true",
                        help="Enable verbose logging")
    args = parser.parse_args()

    logging.basicConfig(
        level=logging.DEBUG if args.verbose else logging.INFO,
        format="%(levelname)s: %(message)s"
    )

    tex_path = Path(args.input).resolve()
    if not tex_path.exists():
        print(f"Error: {tex_path} not found")
        sys.exit(1)

    tex_dir = tex_path.parent
    img_dir = (tex_dir / ".." / "QAL_Paper" / "images").resolve()
    cache_dir = tex_dir / "_pptx_cache"
    pdf_path = Path(args.pdf) if args.pdf else tex_path.with_suffix(".pdf")
    output_path = Path(args.output) if args.output else tex_path.with_name(
        tex_path.stem + "_parsed.pptx"
    )

    # Step 1: Convert PDF figures
    print("Converting PDF figures to PNG...")
    if img_dir.exists():
        convert_pdf_figures(img_dir, cache_dir)
    else:
        print(f"  Warning: Image directory not found: {img_dir}")

    # Step 2: Parse .tex file
    print(f"\nParsing {tex_path.name}...")
    source = tex_path.read_text(encoding="utf-8")
    tex_parser = BeamerParser(source)
    preamble, frames = tex_parser.parse()

    print(f"  Colors defined: {list(preamble.colors.keys())}")
    print(f"  Title: {preamble.title}")
    print(f"  Frames: {len(frames)}")

    # Step 3: Layout
    print("\nComputing layout...")
    engine = LayoutEngine(preamble)
    all_layouts = [engine.layout_frame(f) for f in frames]

    # Step 4: Render PPTX
    print("Generating PPTX...")
    renderer = PptxRenderer(preamble, tex_dir, img_dir, pdf_path, cache_dir)
    prs = renderer.render_presentation(frames, all_layouts)

    # Step 5: Save
    prs.save(str(output_path))
    print(f"\nSaved: {output_path}")
    print(f"  {len(frames)} slides, {LayoutEngine.SLIDE_W}x{LayoutEngine.SLIDE_H} in (16:9)")

    if not pdf_path.exists():
        print(f"\n  Note: {pdf_path} not found - complex TikZ diagrams will show placeholders.")
        print(f"  Run: cd {tex_dir} && pdflatex {tex_path.name}")


if __name__ == "__main__":
    main()
