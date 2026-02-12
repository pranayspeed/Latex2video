#!/usr/bin/env python3
import asyncio
import re
import shutil
import subprocess
import tempfile
import threading
import textwrap
import uuid
from pathlib import Path

from flask import Flask, jsonify, render_template, request, send_file
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from werkzeug.utils import secure_filename

from moviepy.editor import AudioFileClip, ImageClip, concatenate_videoclips
from PIL import Image, ImageDraw, ImageFont
import edge_tts

APP_ROOT = Path(__file__).resolve().parent
ALLOWED_EXT = {".pptx"}
MAX_CONTENT_LENGTH = 300 * 1024 * 1024
ALLOWED_VIDEO_RATES = {"-20%", "-10%", "+0%", "+10%", "+20%", "+30%", "+40%", "+50%"}

JOBS: dict[str, dict] = {}
JOBS_LOCK = threading.Lock()

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = MAX_CONTENT_LENGTH


class VideoBuildError(Exception):
    pass


def _append_log(job_id: str, message: str) -> None:
    with JOBS_LOCK:
        if job_id in JOBS:
            JOBS[job_id]["log"].append(message)


def _set_job_state(job_id: str, **kwargs) -> None:
    with JOBS_LOCK:
        if job_id in JOBS:
            JOBS[job_id].update(kwargs)


def _which_or_raise(*bins: str) -> str:
    for b in bins:
        p = shutil.which(b)
        if p:
            return p
    raise VideoBuildError(f"Required tool not found. Tried: {', '.join(bins)}")


def _run_cmd(cmd: list[str], cwd: Path, log_cb=None) -> None:
    if log_cb:
        log_cb(f"$ {' '.join(cmd)}")
    proc = subprocess.Popen(
        cmd,
        cwd=cwd,
        stdout=subprocess.PIPE,
        stderr=subprocess.STDOUT,
        text=True,
        bufsize=1,
    )
    assert proc.stdout is not None
    for line in proc.stdout:
        line = line.rstrip()
        if line and log_cb:
            log_cb(line)
    rc = proc.wait()
    if rc != 0:
        raise subprocess.CalledProcessError(rc, cmd)


def _clean_note_text(text: str) -> str:
    text = text.replace("\r", "\n")
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _extract_notes_from_pptx(pptx_path: Path) -> list[str]:
    prs = Presentation(str(pptx_path))
    notes: list[str] = []
    for slide in prs.slides:
        note_text = ""
        try:
            notes_slide = slide.notes_slide
            notes_frame = notes_slide.notes_text_frame
            if notes_frame:
                note_text = notes_frame.text or ""
        except Exception:
            note_text = ""

        note_text = _clean_note_text(note_text)
        if note_text.lower() == "click to add notes":
            note_text = ""
        notes.append(note_text)
    return notes


def _pptx_to_pdf(pptx_path: Path, out_pdf: Path, workdir: Path, log_cb=None) -> None:
    soffice = _which_or_raise("soffice", "libreoffice")
    cmd = [
        soffice,
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(workdir),
        str(pptx_path),
    ]
    _run_cmd(cmd, cwd=workdir, log_cb=log_cb)
    produced_pdf = workdir / f"{pptx_path.stem}.pdf"
    if not produced_pdf.exists():
        raise VideoBuildError("PPTX to PDF conversion finished but PDF not found.")
    shutil.move(str(produced_pdf), str(out_pdf))


def _pdf_to_pngs(pdf_path: Path, out_dir: Path, dpi: int, log_cb=None) -> list[Path]:
    _which_or_raise("pdftoppm")
    out_dir.mkdir(parents=True, exist_ok=True)
    prefix = out_dir / "slide"
    cmd = ["pdftoppm", "-r", str(dpi), "-png", str(pdf_path), str(prefix)]
    _run_cmd(cmd, cwd=out_dir.parent, log_cb=log_cb)
    images = sorted(out_dir.glob("slide-*.png"))
    if not images:
        raise VideoBuildError("No slide images generated from PDF.")
    return images


def _extract_slide_images_from_pptx(pptx_path: Path, out_dir: Path, log_cb=None) -> list[Path]:
    prs = Presentation(str(pptx_path))
    out_dir.mkdir(parents=True, exist_ok=True)
    images: list[Path] = []

    for idx, slide in enumerate(prs.slides, start=1):
        best_shape = None
        best_area = 0
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            area = int(shape.width) * int(shape.height)
            if area > best_area:
                best_area = area
                best_shape = shape

        if best_shape is None:
            raise VideoBuildError(
                "Slide image extraction fallback failed: no picture found on one or more slides. "
                "Install LibreOffice (soffice) for full PPTX rendering support."
            )

        ext = best_shape.image.ext or "png"
        out_path = out_dir / f"slide-{idx:03d}.{ext}"
        out_path.write_bytes(best_shape.image.blob)
        images.append(out_path)

    if log_cb:
        log_cb(f"Extracted {len(images)} slide images directly from PPTX")
    return images


def _load_font(size: int, bold: bool = False) -> ImageFont.ImageFont:
    candidates = [
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/Library/Fonts/Arial.ttf",
        "/Library/Fonts/Arial Bold.ttf" if bold else "/Library/Fonts/Arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    for fp in candidates:
        try:
            return ImageFont.truetype(fp, size=size)
        except Exception:
            continue
    return ImageFont.load_default()


def _draw_wrapped_text(
    draw: ImageDraw.ImageDraw,
    text: str,
    font: ImageFont.ImageFont,
    x: int,
    y: int,
    max_width: int,
    fill: tuple[int, int, int],
    line_spacing: int,
) -> int:
    if not text.strip():
        return y

    avg_char_px = max(7, int(font.size * 0.55)) if hasattr(font, "size") else 10
    wrap_width = max(20, max_width // avg_char_px)
    lines = textwrap.wrap(text, width=wrap_width, break_long_words=False, break_on_hyphens=False) or [text]

    for line in lines:
        draw.text((x, y), line, font=font, fill=fill)
        bbox = draw.textbbox((x, y), line, font=font)
        line_h = max(18, bbox[3] - bbox[1])
        y += line_h + line_spacing
    return y


def _extract_slide_text(slide, idx: int) -> tuple[str, list[str]]:
    title = ""
    try:
        if slide.shapes.title and getattr(slide.shapes.title, "has_text_frame", False):
            title = _clean_note_text(slide.shapes.title.text_frame.text or "")
    except Exception:
        title = ""

    body_chunks: list[str] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        txt = _clean_note_text(shape.text_frame.text or "")
        if not txt:
            continue
        if txt.lower() in {"click to add text", "click to add title"}:
            continue
        if title and txt == title:
            continue
        body_chunks.append(txt)

    if not title:
        title = f"Slide {idx}"
    return title, body_chunks


def _render_text_slides_from_pptx(pptx_path: Path, out_dir: Path, log_cb=None) -> list[Path]:
    prs = Presentation(str(pptx_path))
    out_dir.mkdir(parents=True, exist_ok=True)

    width_ratio = float(prs.slide_width) if prs.slide_width else 16.0
    height_ratio = float(prs.slide_height) if prs.slide_height else 9.0
    out_w = 1280
    out_h = max(720, int(out_w * (height_ratio / width_ratio)))

    title_font = _load_font(54, bold=True)
    body_font = _load_font(34, bold=False)
    images: list[Path] = []

    for idx, slide in enumerate(prs.slides, start=1):
        title, body_chunks = _extract_slide_text(slide, idx)
        img = Image.new("RGB", (out_w, out_h), color=(248, 250, 252))
        draw = ImageDraw.Draw(img)

        pad_x = int(out_w * 0.08)
        max_w = int(out_w * 0.84)
        y = int(out_h * 0.10)

        y = _draw_wrapped_text(
            draw=draw,
            text=title,
            font=title_font,
            x=pad_x,
            y=y,
            max_width=max_w,
            fill=(15, 23, 42),
            line_spacing=8,
        )
        y += 20

        if body_chunks:
            for chunk in body_chunks:
                y = _draw_wrapped_text(
                    draw=draw,
                    text=chunk,
                    font=body_font,
                    x=pad_x,
                    y=y,
                    max_width=max_w,
                    fill=(30, 41, 59),
                    line_spacing=6,
                )
                y += 16
                if y > int(out_h * 0.88):
                    break
        else:
            _draw_wrapped_text(
                draw=draw,
                text="(No slide text found in PPTX content)",
                font=body_font,
                x=pad_x,
                y=y,
                max_width=max_w,
                fill=(71, 85, 105),
                line_spacing=6,
            )

        out_path = out_dir / f"slide-{idx:03d}.png"
        img.save(out_path, format="PNG")
        images.append(out_path)

    if log_cb:
        log_cb(f"Rendered {len(images)} text-based slide images from PPTX")
    return images


async def _tts_to_file(text: str, output_path: Path, voice: str, rate: str) -> None:
    communicator = edge_tts.Communicate(text, voice=voice, rate=rate)
    await communicator.save(str(output_path))


def _generate_audio(text: str, output_path: Path, voice: str, rate: str) -> None:
    asyncio.run(_tts_to_file(text, output_path, voice=voice, rate=rate))


def _build_video(
    images: list[Path],
    notes: list[str],
    out_video: Path,
    voice: str,
    rate: str,
    fps: int,
    default_seconds: float,
    workdir: Path,
    log_cb=None,
) -> None:
    _which_or_raise("ffmpeg")
    audio_dir = workdir / "audio"
    audio_dir.mkdir(parents=True, exist_ok=True)

    pair_count = min(len(images), len(notes))
    if pair_count == 0:
        raise VideoBuildError("No slides found.")

    clips = []
    for idx in range(pair_count):
        img_path = images[idx]
        note = (notes[idx] or "").strip()
        if log_cb:
            log_cb(f"Preparing slide {idx + 1}/{pair_count}")
        img_clip = ImageClip(str(img_path))

        if note:
            audio_path = audio_dir / f"slide_{idx + 1}.mp3"
            _generate_audio(note, audio_path, voice=voice, rate=rate)
            audio_clip = AudioFileClip(str(audio_path))
            clip = img_clip.set_duration(audio_clip.duration).set_audio(audio_clip)
        else:
            clip = img_clip.set_duration(default_seconds)

        clips.append(clip)

    final = concatenate_videoclips(clips, method="compose")
    if log_cb:
        log_cb(f"Rendering final video ({final.duration:.1f}s)")
    final.write_videofile(str(out_video), fps=fps, codec="libx264", audio_codec="aac")

    final.close()
    for c in clips:
        c.close()


def _process_job(job_id: str, filename: str, payload: bytes, video_rate: str) -> None:
    tmp_dir = Path(tempfile.mkdtemp(prefix="pptx_video_job_"))
    _set_job_state(job_id, temp_dir=str(tmp_dir))
    log = lambda m: _append_log(job_id, m)

    try:
        pptx_path = tmp_dir / filename
        pptx_path.write_bytes(payload)
        log(f"Saved upload: {filename}")

        notes = _extract_notes_from_pptx(pptx_path)
        log(f"Extracted notes from {len(notes)} slides")

        images = []
        try:
            out_pdf = tmp_dir / "slides.pdf"
            _pptx_to_pdf(pptx_path, out_pdf, tmp_dir, log_cb=log)
            log("Converted PPTX to PDF")
            images = _pdf_to_pngs(out_pdf, tmp_dir / "slides", dpi=220, log_cb=log)
            log(f"Rendered {len(images)} slide images")
        except VideoBuildError as e:
            if "Required tool not found" in str(e):
                log("LibreOffice not found, using direct PPTX image extraction fallback.")
                try:
                    images = _extract_slide_images_from_pptx(pptx_path, tmp_dir / "slides_from_pptx", log_cb=log)
                except VideoBuildError:
                    log("No embedded slide pictures found. Falling back to text-rendered slide images.")
                    images = _render_text_slides_from_pptx(pptx_path, tmp_dir / "slides_text_fallback", log_cb=log)
            else:
                raise

        out_video = tmp_dir / f"{pptx_path.stem}_narrated.mp4"
        _build_video(
            images=images,
            notes=notes,
            out_video=out_video,
            voice="en-US-ChristopherNeural",
            rate=video_rate,
            fps=24,
            default_seconds=3.0,
            workdir=tmp_dir,
            log_cb=log,
        )
        _set_job_state(job_id, status="done", output_path=str(out_video), output_name=out_video.name)
        log("Video generated.")
    except subprocess.CalledProcessError as e:
        _set_job_state(job_id, status="error", error=f"Command failed: {' '.join(e.cmd)}")
        log(f"Error: command failed with exit code {e.returncode}")
    except Exception as e:
        _set_job_state(job_id, status="error", error=str(e))
        log(f"Error: {e}")


@app.route("/")
def index():
    return render_template("pptx_video.html")


@app.route("/convert", methods=["POST"])
def convert():
    if "file" not in request.files:
        return jsonify({"error": "Missing file upload."}), 400

    f = request.files["file"]
    if f.filename == "":
        return jsonify({"error": "No file selected."}), 400

    filename = secure_filename(f.filename)
    ext = Path(filename).suffix.lower()
    if ext not in ALLOWED_EXT:
        return jsonify({"error": "Only .pptx files are supported."}), 400

    video_rate = request.form.get("video_rate", "+20%")
    if video_rate not in ALLOWED_VIDEO_RATES:
        return jsonify({"error": "Invalid narration speed."}), 400

    payload = f.read()
    job_id = uuid.uuid4().hex
    with JOBS_LOCK:
        JOBS[job_id] = {
            "status": "running",
            "log": [f"Started job {job_id} (pptx_video, rate={video_rate})"],
            "error": None,
            "output_path": None,
            "output_name": None,
            "temp_dir": None,
        }

    thread = threading.Thread(target=_process_job, args=(job_id, filename, payload, video_rate), daemon=True)
    thread.start()
    return jsonify({"job_id": job_id})


@app.route("/status/<job_id>", methods=["GET"])
def status(job_id: str):
    with JOBS_LOCK:
        job = JOBS.get(job_id)
        if not job:
            return jsonify({"error": "Job not found."}), 404
        return jsonify(
            {
                "status": job["status"],
                "log": "\n".join(job["log"]),
                "error": job["error"],
                "ready": bool(job["output_path"]),
            }
        )


@app.route("/download/<job_id>", methods=["GET"])
def download(job_id: str):
    with JOBS_LOCK:
        job = JOBS.get(job_id)
        if not job:
            return jsonify({"error": "Job not found."}), 404
        output_path = job["output_path"]
        output_name = job["output_name"]

    if not output_path or not output_name:
        return jsonify({"error": "Output not ready."}), 409

    output_file = Path(output_path)
    if not output_file.exists():
        return jsonify({"error": "Output file missing."}), 410

    return send_file(output_file, as_attachment=True, download_name=output_name)


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=8001, debug=True)
